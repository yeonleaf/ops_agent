"""적응형 파일 프로세서 - 페이지별로 최적의 처리 방식 선택"""

import os
import tempfile
import json
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd

from module.image_to_text import AzureOpenAIImageProcessor
from module.types import DocumentType, ContentType, ElementType, ClassificationResult
from module.page_analyzer import PageContentAnalyzer
from module.file_processor import DocumentElement, ProcessedPage, TextBasedProcessor, LayoutBasedProcessor


class AdaptiveFileProcessor:
    """페이지별로 최적의 처리 방식을 선택하는 적응형 파일 프로세서"""
    
    def __init__(self, azure_processor: AzureOpenAIImageProcessor):
        self.azure_processor = azure_processor
        self.page_analyzer = PageContentAnalyzer()
        self.text_processor = TextBasedProcessor(azure_processor)
        self.layout_processor = LayoutBasedProcessor(azure_processor)
        self.temp_storage = {}
        
        # 처리 통계
        self.processing_stats = {
            "text_based_pages": 0,
            "layout_based_pages": 0,
            "total_pages": 0,
            "page_classifications": []
        }
    
    def process_file_adaptive(self, file_path: str, doc_type: DocumentType) -> Dict[str, Any]:
        """파일을 페이지별 적응형 방식으로 처리"""
        try:
            print(f"\n🚀 적응형 처리 시작: {file_path}")
            print(f"📄 파일 타입: {doc_type.value}")
            
            # 페이지별 처리
            if doc_type == DocumentType.PDF:
                processed_pages = self._process_pdf_adaptive(file_path)
            elif doc_type == DocumentType.DOCX:
                processed_pages = self._process_docx_adaptive(file_path)
            elif doc_type == DocumentType.PPTX:
                processed_pages = self._process_pptx_adaptive(file_path)
            elif doc_type == DocumentType.XLSX:
                processed_pages = self._process_xlsx_adaptive(file_path)
            elif doc_type == DocumentType.XML:
                processed_pages = self.text_processor._process_xml_file(file_path)
            else:
                raise ValueError(f"지원하지 않는 파일 타입: {doc_type}")
            
            # 전체 결과 생성
            result = {
                "file_path": file_path,
                "file_type": doc_type.value,
                "processing_mode": "adaptive",
                "processed_pages": [page.to_dict() for page in processed_pages],
                "total_pages": len(processed_pages),
                "processing_stats": self.processing_stats,
                "processing_timestamp": str(pd.Timestamp.now())
            }
            
            # 처리 통계 출력
            self._print_processing_summary()
            
            # 임시 저장
            self._save_to_temp_storage(file_path, result)
            
            return result
            
        except Exception as e:
            print(f"❌ 적응형 처리 오류: {e}")
            return {"error": str(e)}
    
    def _process_pdf_adaptive(self, file_path: str) -> List[ProcessedPage]:
        """PDF를 페이지별 적응형으로 처리"""
        try:
            doc = fitz.open(file_path)
            processed_pages = []
            
            for page_num in range(1, len(doc) + 1):
                print(f"📄 페이지 {page_num}/{len(doc)} 분석 중...")
                
                # 페이지별 콘텐츠 분석
                classification = self.page_analyzer.analyze_pdf_page(file_path, page_num)
                self._update_stats(page_num, classification)
                
                # 분류 결과에 따른 처리
                if classification.content_type == ContentType.TEXT_BASED:
                    page = self._process_pdf_page_text_based(doc, page_num - 1, classification)
                else:
                    page = self._process_pdf_page_layout_based(file_path, page_num, classification)
                
                processed_pages.append(page)
            
            doc.close()
            return processed_pages
            
        except Exception as e:
            print(f"PDF 적응형 처리 오류: {e}")
            return []
    
    def _process_docx_adaptive(self, file_path: str) -> List[ProcessedPage]:
        """DOCX를 페이지별 적응형으로 처리"""
        try:
            doc = Document(file_path)
            processed_pages = []
            
            # 문단을 페이지로 분할 (대략적)
            paragraphs_per_page = 12
            total_paragraphs = len(doc.paragraphs)
            total_pages = (total_paragraphs + paragraphs_per_page - 1) // paragraphs_per_page
            
            for page_num in range(1, total_pages + 1):
                print(f"📄 페이지 {page_num}/{total_pages} 분석 중...")
                
                # 페이지별 콘텐츠 분석
                classification = self.page_analyzer.analyze_docx_page(file_path, page_num)
                self._update_stats(page_num, classification)
                
                # 분류 결과에 따른 처리
                if classification.content_type == ContentType.TEXT_BASED:
                    page = self._process_docx_page_text_based(doc, page_num, paragraphs_per_page, classification)
                else:
                    page = self._process_docx_page_layout_based(file_path, page_num, classification)
                
                processed_pages.append(page)
            
            return processed_pages
            
        except Exception as e:
            print(f"DOCX 적응형 처리 오류: {e}")
            return []
    
    def _process_pptx_adaptive(self, file_path: str) -> List[ProcessedPage]:
        """PPTX를 슬라이드별 적응형으로 처리"""
        try:
            prs = Presentation(file_path)
            processed_pages = []
            
            for slide_num in range(1, len(prs.slides) + 1):
                print(f"📊 슬라이드 {slide_num}/{len(prs.slides)} 분석 중...")
                
                # 슬라이드별 콘텐츠 분석
                classification = self.page_analyzer.analyze_pptx_slide(file_path, slide_num)
                self._update_stats(slide_num, classification)
                
                # 분류 결과에 따른 처리
                if classification.content_type == ContentType.TEXT_BASED:
                    page = self._process_pptx_slide_text_based(prs, slide_num - 1, classification)
                else:
                    page = self._process_pptx_slide_layout_based(file_path, slide_num, classification)
                
                processed_pages.append(page)
            
            return processed_pages
            
        except Exception as e:
            print(f"PPTX 적응형 처리 오류: {e}")
            return []
    
    def _process_xlsx_adaptive(self, file_path: str) -> List[ProcessedPage]:
        """XLSX를 시트별 적응형으로 처리"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            processed_pages = []
            
            for sheet_num in range(1, len(wb.sheetnames) + 1):
                print(f"📋 시트 {sheet_num}/{len(wb.sheetnames)} 분석 중...")
                
                # 시트별 콘텐츠 분석
                classification = self.page_analyzer.analyze_xlsx_sheet(file_path, sheet_num)
                self._update_stats(sheet_num, classification)
                
                # Excel은 대부분 text-based로 처리
                page = self._process_xlsx_sheet_text_based(wb, sheet_num, classification)
                processed_pages.append(page)
            
            wb.close()
            return processed_pages
            
        except Exception as e:
            print(f"XLSX 적응형 처리 오류: {e}")
            return []
    
    def _process_pdf_page_text_based(self, doc: fitz.Document, page_idx: int, 
                                    classification: ClassificationResult) -> ProcessedPage:
        """PDF 페이지를 텍스트 기반으로 처리"""
        try:
            page = doc[page_idx]
            text = page.get_text()
            
            elements = [DocumentElement(
                ElementType.TEXT,
                text,
                {
                    "page_number": page_idx + 1,
                    "source": "pdf_text_extraction",
                    "classification": classification.content_type.value,
                    "confidence": classification.confidence,
                    "reasoning": classification.reasoning[:3]
                }
            )]
            
            return ProcessedPage(
                page_idx + 1, elements, "page",
                {
                    "file_type": "pdf",
                    "processing_method": "text_based",
                    "classification": classification.content_type.value,
                    "confidence": classification.confidence
                }
            )
            
        except Exception as e:
            print(f"PDF 페이지 {page_idx + 1} 텍스트 처리 오류: {e}")
            return ProcessedPage(page_idx + 1, [], "page", {"error": str(e)})
    
    def _process_pdf_page_layout_based(self, file_path: str, page_num: int, 
                                      classification: ClassificationResult) -> ProcessedPage:
        """PDF 페이지를 레이아웃 기반으로 처리 (이미지 변환 후 GPT Vision)"""
        try:
            # 단일 페이지를 이미지로 변환
            doc = fitz.open(file_path)
            page = doc[page_num - 1]
            
            # 이미지로 변환
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x 해상도
            img_data = pix.tobytes("png")
            
            # GPT Vision으로 처리
            try:
                vision_result = self.azure_processor.process_image_from_bytes(
                    img_data, 
                    f"페이지 {page_num}의 모든 텍스트를 추출하고 표나 이미지가 있다면 설명해주세요."
                )
                
                elements = [DocumentElement(
                    ElementType.TEXT,
                    vision_result,
                    {
                        "page_number": page_num,
                        "source": "gpt_vision",
                        "classification": classification.content_type.value,
                        "confidence": classification.confidence,
                        "reasoning": classification.reasoning[:3]
                    }
                )]
                
            except Exception as vision_error:
                print(f"GPT Vision 처리 오류: {vision_error}")
                # 폴백: 텍스트 추출
                text = page.get_text()
                elements = [DocumentElement(
                    ElementType.TEXT,
                    text or "텍스트 추출 실패",
                    {
                        "page_number": page_num,
                        "source": "pdf_text_fallback",
                        "error": str(vision_error)
                    }
                )]
            
            doc.close()
            
            return ProcessedPage(
                page_num, elements, "page",
                {
                    "file_type": "pdf",
                    "processing_method": "layout_based",
                    "classification": classification.content_type.value,
                    "confidence": classification.confidence
                }
            )
            
        except Exception as e:
            print(f"PDF 페이지 {page_num} 레이아웃 처리 오류: {e}")
            return ProcessedPage(page_num, [], "page", {"error": str(e)})
    
    def _process_docx_page_text_based(self, doc: Document, page_num: int, paragraphs_per_page: int,
                                     classification: ClassificationResult) -> ProcessedPage:
        """DOCX 페이지를 텍스트 기반으로 처리"""
        try:
            start_idx = (page_num - 1) * paragraphs_per_page
            end_idx = min(start_idx + paragraphs_per_page, len(doc.paragraphs))
            
            page_paragraphs = doc.paragraphs[start_idx:end_idx]
            text_content = "\n".join([p.text for p in page_paragraphs if p.text.strip()])
            
            elements = [DocumentElement(
                ElementType.TEXT,
                text_content,
                {
                    "page_number": page_num,
                    "source": "docx_text_extraction",
                    "classification": classification.content_type.value,
                    "confidence": classification.confidence,
                    "paragraph_range": f"{start_idx}-{end_idx}"
                }
            )]
            
            return ProcessedPage(
                page_num, elements, "page",
                {
                    "file_type": "docx",
                    "processing_method": "text_based",
                    "classification": classification.content_type.value
                }
            )
            
        except Exception as e:
            print(f"DOCX 페이지 {page_num} 텍스트 처리 오류: {e}")
            return ProcessedPage(page_num, [], "page", {"error": str(e)})
    
    def _process_docx_page_layout_based(self, file_path: str, page_num: int,
                                       classification: ClassificationResult) -> ProcessedPage:
        """DOCX 페이지를 레이아웃 기반으로 처리 (PDF 변환 후 이미지 처리)"""
        try:
            # DOCX를 PDF로 변환 후 해당 페이지만 이미지로 처리
            # 임시 PDF 생성
            temp_dir = "temp_pdfs"
            os.makedirs(temp_dir, exist_ok=True)
            pdf_path = os.path.join(temp_dir, f"temp_docx_page_{page_num}.pdf")
            
            # LibreOffice로 변환 (기존 레이아웃 프로세서 활용)
            self.layout_processor._convert_docx_to_pdf(file_path, pdf_path)
            
            # PDF의 해당 페이지를 이미지로 처리
            result_page = self._process_pdf_page_layout_based(pdf_path, page_num, classification)
            
            # 메타데이터 업데이트
            result_page.metadata["original_file_type"] = "docx"
            result_page.metadata["converted_via"] = "pdf"
            
            return result_page
            
        except Exception as e:
            print(f"DOCX 페이지 {page_num} 레이아웃 처리 오류: {e}")
            # 폴백: 텍스트 기반 처리
            return self._process_docx_page_text_based(Document(file_path), page_num, 12, classification)
    
    def _process_pptx_slide_text_based(self, prs: Presentation, slide_idx: int,
                                      classification: ClassificationResult) -> ProcessedPage:
        """PPTX 슬라이드를 텍스트 기반으로 처리"""
        try:
            slide = prs.slides[slide_idx]
            text_content = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text.strip() + "\n"
            
            elements = [DocumentElement(
                ElementType.TEXT,
                text_content,
                {
                    "page_number": slide_idx + 1,
                    "source": "pptx_text_extraction",
                    "classification": classification.content_type.value,
                    "confidence": classification.confidence
                }
            )]
            
            return ProcessedPage(
                slide_idx + 1, elements, "slide",
                {
                    "file_type": "pptx",
                    "processing_method": "text_based",
                    "classification": classification.content_type.value
                }
            )
            
        except Exception as e:
            print(f"PPTX 슬라이드 {slide_idx + 1} 텍스트 처리 오류: {e}")
            return ProcessedPage(slide_idx + 1, [], "slide", {"error": str(e)})
    
    def _process_pptx_slide_layout_based(self, file_path: str, slide_num: int,
                                        classification: ClassificationResult) -> ProcessedPage:
        """PPTX 슬라이드를 레이아웃 기반으로 처리"""
        try:
            # PPTX를 PDF로 변환 후 해당 슬라이드만 이미지로 처리
            temp_dir = "temp_pdfs"
            os.makedirs(temp_dir, exist_ok=True)
            pdf_path = os.path.join(temp_dir, f"temp_pptx_slide_{slide_num}.pdf")
            
            # PPTX를 PDF로 변환
            self.layout_processor._convert_pptx_to_pdf(file_path, pdf_path)
            
            # PDF의 해당 슬라이드를 이미지로 처리
            result_page = self._process_pdf_page_layout_based(pdf_path, slide_num, classification)
            
            # 메타데이터 업데이트
            result_page.metadata["original_file_type"] = "pptx"
            result_page.metadata["converted_via"] = "pdf"
            result_page.page_type = "slide"
            
            return result_page
            
        except Exception as e:
            print(f"PPTX 슬라이드 {slide_num} 레이아웃 처리 오류: {e}")
            # 폴백: 텍스트 기반 처리
            return self._process_pptx_slide_text_based(Presentation(file_path), slide_num - 1, classification)
    
    def _process_xlsx_sheet_text_based(self, wb, sheet_num: int,
                                      classification: ClassificationResult) -> ProcessedPage:
        """XLSX 시트를 텍스트 기반으로 처리"""
        try:
            sheet_name = wb.sheetnames[sheet_num - 1]
            ws = wb[sheet_name]
            
            # 시트 데이터를 표로 변환
            table_data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    table_data.append([str(cell) if cell is not None else "" for cell in row])
            
            elements = [DocumentElement(
                ElementType.TABLE,
                table_data,
                {
                    "page_number": sheet_num,
                    "source": "xlsx_direct_extraction",
                    "classification": classification.content_type.value,
                    "sheet_name": sheet_name
                }
            )]
            
            return ProcessedPage(
                sheet_num, elements, "sheet",
                {
                    "file_type": "xlsx",
                    "processing_method": "text_based",
                    "sheet_name": sheet_name,
                    "classification": classification.content_type.value
                }
            )
            
        except Exception as e:
            print(f"XLSX 시트 {sheet_num} 처리 오류: {e}")
            return ProcessedPage(sheet_num, [], "sheet", {"error": str(e)})
    
    def _update_stats(self, page_num: int, classification: ClassificationResult):
        """처리 통계 업데이트"""
        self.processing_stats["total_pages"] += 1
        
        if classification.content_type == ContentType.TEXT_BASED:
            self.processing_stats["text_based_pages"] += 1
        else:
            self.processing_stats["layout_based_pages"] += 1
        
        self.processing_stats["page_classifications"].append({
            "page_number": page_num,
            "content_type": classification.content_type.value,
            "confidence": classification.confidence,
            "reasoning": classification.reasoning[:2]  # 상위 2개만
        })
        
        # 실시간 진행상황 출력
        content_type_icon = "📝" if classification.content_type == ContentType.TEXT_BASED else "🎨"
        print(f"   {content_type_icon} {classification.content_type.value} (신뢰도: {classification.confidence:.2f})")
    
    def _print_processing_summary(self):
        """처리 결과 요약 출력"""
        stats = self.processing_stats
        total = stats["total_pages"]
        text_count = stats["text_based_pages"]
        layout_count = stats["layout_based_pages"]
        
        print(f"\n📊 적응형 처리 완료 요약")
        print(f"   총 페이지: {total}개")
        if total > 0:
            print(f"   📝 텍스트 기반: {text_count}개 ({text_count/total*100:.1f}%)")
            print(f"   🎨 레이아웃 기반: {layout_count}개 ({layout_count/total*100:.1f}%)")
        else:
            print(f"   📝 텍스트 기반: {text_count}개")
            print(f"   🎨 레이아웃 기반: {layout_count}개")
        
        # 페이지별 분류 결과 간략 표시 (처음 10개만)
        print(f"\n📄 페이지별 분류 결과 (처음 10개):")
        for i, classification in enumerate(stats["page_classifications"][:10]):
            page_num = classification["page_number"]
            content_type = classification["content_type"]
            confidence = classification["confidence"]
            icon = "📝" if content_type == "text_based" else "🎨"
            print(f"   페이지 {page_num:2d}: {icon} {content_type} ({confidence:.2f})")
        
        if total > 10:
            print(f"   ... 및 {total - 10}개 페이지 더")
    
    def _save_to_temp_storage(self, file_path: str, result: Dict[str, Any]):
        """결과를 임시 저장소에 저장"""
        try:
            temp_dir = "temp_results"
            os.makedirs(temp_dir, exist_ok=True)
            
            filename = Path(file_path).stem
            result_filename = f"{filename}_adaptive_result.json"
            result_path = os.path.join(temp_dir, result_filename)
            
            with open(result_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            
            print(f"📁 적응형 처리 결과 저장: {result_path}")
            
        except Exception as e:
            print(f"결과 저장 오류: {e}")