import os
import sys
import json
import tempfile
import shutil
from typing import List, Dict, Any, Optional, Tuple, Union
from enum import Enum
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook
import xlrd
from PIL import Image
import io
import xml.etree.ElementTree as ET
import html

from module.image_to_text import AzureOpenAIImageProcessor
from structured_chunking import JiraStructuredChunker, StructuredChunk

class DocumentType(Enum):
    """문서 타입"""
    DOCX = "docx"
    PPTX = "pptx"
    PDF = "pdf"
    XLSX = "xlsx"
    XLS = "xls"
    HTML = "html"
    TEXT = "text"
    XML = "xml"

class ContentType(Enum):
    """콘텐츠 타입 (text-based vs layout-based)"""
    TEXT_BASED = "text_based"
    LAYOUT_BASED = "layout_based"

class ElementType(Enum):
    """문서 요소 타입"""
    TEXT = "text"
    TABLE = "table"
    IMAGE = "image"
    SHAPE = "shape"

class DocumentElement:
    """문서 요소 클래스"""
    def __init__(self, element_type: ElementType, content: Any, metadata: Dict = None):
        self.element_type = element_type
        self.content = content
        self.metadata = metadata or {}
    
    def __str__(self):
        return f"{self.element_type.value}: {str(self.content)[:50]}..."

    def to_dict(self):
        """딕셔너리로 변환"""
        return {
            "element_type": self.element_type.value,
            "content": self.content,
            "metadata": self.metadata
        }

class ProcessedPage:
    """처리된 페이지/슬라이드 정보"""
    def __init__(self, page_number: int, elements: List[DocumentElement], 
                 page_type: str = "page", metadata: Dict = None):
        self.page_number = page_number
        self.elements = elements
        self.page_type = page_type  # "page", "slide", "sheet"
        self.metadata = metadata or {}
    
    def to_dict(self):
        """딕셔너리로 변환"""
        return {
            "page_number": self.page_number,
            "page_type": self.page_type,
            "elements": [elem.to_dict() for elem in self.elements],
            "metadata": self.metadata
        }

class FileTypeDetector:
    """파일 타입 판별기"""
    
    @staticmethod
    def detect_file_type(file_path: str) -> DocumentType:
        """파일 확장자로부터 문서 타입 판별"""
        ext = Path(file_path).suffix.lower()
        if ext == '.docx':
            return DocumentType.DOCX
        elif ext == '.pptx':
            return DocumentType.PPTX
        elif ext == '.pdf':
            return DocumentType.PDF
        elif ext == '.xlsx':
            return DocumentType.XLSX
        elif ext == '.xls':
            return DocumentType.XLS
        elif ext == '.html':
            return DocumentType.HTML
        elif ext in ['.txt', '.md']:
            return DocumentType.TEXT
        elif ext == '.xml':
            return DocumentType.XML
        else:
            raise ValueError(f"지원하지 않는 파일 타입: {ext}")
    
    @staticmethod
    def detect_content_type(file_path: str, doc_type: DocumentType) -> ContentType:
        """콘텐츠가 text-based인지 layout-based인지 판별"""
        try:
            if doc_type == DocumentType.PDF:
                return FileTypeDetector._analyze_pdf_content(file_path)
            elif doc_type == DocumentType.DOCX:
                return FileTypeDetector._analyze_docx_content(file_path)
            elif doc_type == DocumentType.PPTX:
                return FileTypeDetector._analyze_pptx_content(file_path)
            elif doc_type == DocumentType.XLSX:
                return FileTypeDetector._analyze_xlsx_content(file_path)
            elif doc_type == DocumentType.XLS:
                return FileTypeDetector._analyze_xls_content(file_path)
            elif doc_type == DocumentType.XML:
                return ContentType.TEXT_BASED  # XML은 항상 text-based로 처리
            else:
                return ContentType.TEXT_BASED
        except Exception as e:
            print(f"콘텐츠 타입 판별 오류: {e}")
            return ContentType.LAYOUT_BASED  # 오류 시 안전하게 layout-based로 분류
    
    @staticmethod
    def _analyze_pdf_content(file_path: str) -> ContentType:
        """PDF 콘텐츠 분석"""
        try:
            doc = fitz.open(file_path)
            text_content = ""
            image_count = 0
            
            for page in doc:
                text_content += page.get_text()
                image_list = page.get_images()
                image_count += len(image_list)
            
            doc.close()
            
            # 텍스트가 충분하고 이미지가 적으면 text-based
            if len(text_content.strip()) > 100 and image_count < 3:
                return ContentType.TEXT_BASED
            else:
                return ContentType.LAYOUT_BASED
                
        except Exception:
            return ContentType.LAYOUT_BASED
    
    @staticmethod
    def _analyze_docx_content(file_path: str) -> ContentType:
        """DOCX 콘텐츠 분석"""
        try:
            doc = Document(file_path)
            text_content = ""
            image_count = 0
            
            for para in doc.paragraphs:
                text_content += para.text + "\n"
            
            # 이미지 개수 확인 (간단한 방법)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
            
            if len(text_content.strip()) > 100 and image_count < 3:
                return ContentType.TEXT_BASED
            else:
                return ContentType.LAYOUT_BASED
                
        except Exception:
            return ContentType.LAYOUT_BASED
    
    @staticmethod
    def _analyze_pptx_content(file_path: str) -> ContentType:
        """PPTX 콘텐츠 분석"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            image_count = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
                    if shape.shape_type == 13:  # 이미지 타입
                        image_count += 1
            
            if len(text_content.strip()) > 50 and image_count < 5:
                return ContentType.TEXT_BASED
            else:
                return ContentType.LAYOUT_BASED
                
        except Exception:
            return ContentType.LAYOUT_BASED
    
    @staticmethod
    def _analyze_xlsx_content(file_path: str) -> ContentType:
        """XLSX 콘텐츠 분석"""
        try:
            wb = load_workbook(file_path, data_only=True)
            text_content = ""
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text_content += str(cell) + " "
            
            wb.close()
            
            # 엑셀은 대부분 text-based
            if len(text_content.strip()) > 50:
                return ContentType.TEXT_BASED
            else:
                return ContentType.LAYOUT_BASED
                
        except Exception:
            return ContentType.LAYOUT_BASED
    
    @staticmethod
    def _analyze_xls_content(file_path: str) -> ContentType:
        """XLS 콘텐츠 분석"""
        try:
            wb = xlrd.open_workbook(file_path)
            text_content = ""
            
            for sheet_name in wb.sheet_names():
                sheet = wb.sheet_by_name(sheet_name)
                for row_idx in range(sheet.nrows):
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        if cell_value:
                            text_content += str(cell_value) + " "
            
            # XLS 파일은 항상 text-based로 처리 (PDF 변환 불가능)
            return ContentType.TEXT_BASED
                
        except Exception:
            # 오류가 발생해도 text-based로 처리 (PDF 변환 방지)
            return ContentType.TEXT_BASED

class TextBasedProcessor:
    """Text-based 파일 처리기"""
    
    def __init__(self, azure_processor: AzureOpenAIImageProcessor):
        self.azure_processor = azure_processor
    
    def process_docx(self, file_path: str) -> List[ProcessedPage]:
        """DOCX 파일을 페이지별로 처리"""
        try:
            doc = Document(file_path)
            elements = []
            
            # 문단별로 요소 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    elements.append(DocumentElement(
                        ElementType.TEXT, 
                        para.text.strip(),
                        {"paragraph_style": para.style.name if para.style else "Normal"}
                    ))
            
            # 테이블 처리
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                if table_data:
                    elements.append(DocumentElement(
                        ElementType.TABLE,
                        table_data,
                        {"table_type": "docx_table"}
                    ))
            
            # 이미지 처리
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_path = rel.target_path
                        elements.append(DocumentElement(
                            ElementType.IMAGE,
                            image_path,
                            {"image_source": "docx_embed"}
                        ))
                    except Exception as e:
                        print(f"이미지 처리 오류: {e}")
            
            # 단일 페이지로 처리 (DOCX는 페이지 구분이 명확하지 않음)
            return [ProcessedPage(1, elements, "page", {"file_type": "docx"})]
            
        except Exception as e:
            print(f"DOCX 처리 오류: {e}")
            return []
    
    def process_pptx(self, file_path: str) -> List[ProcessedPage]:
        """PPTX 파일을 슬라이드별로 처리"""
        try:
            prs = Presentation(file_path)
            processed_slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                elements = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        elements.append(DocumentElement(
                            ElementType.TEXT,
                            shape.text.strip(),
                            {"shape_type": str(shape.shape_type)}
                        ))
                    
                    if shape.shape_type == 13:  # 이미지
                        try:
                            # 이미지 추출 및 저장
                            image_path = self._extract_pptx_image(shape, slide_num, len(elements))
                            elements.append(DocumentElement(
                                ElementType.IMAGE,
                                image_path,
                                {"image_source": "pptx_slide"}
                            ))
                        except Exception as e:
                            print(f"PPTX 이미지 추출 오류: {e}")
                
                processed_slides.append(ProcessedPage(
                    slide_num, elements, "slide", 
                    {"file_type": "pptx", "slide_layout": slide.slide_layout.name}
                ))
            
            return processed_slides
            
        except Exception as e:
            print(f"PPTX 처리 오류: {e}")
            return []
    
    def process_xlsx(self, file_path: str) -> List[ProcessedPage]:
        """XLSX 파일을 시트별로 처리"""
        try:
            wb = load_workbook(file_path, data_only=True)
            processed_sheets = []
            
            for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
                ws = wb[sheet_name]
                elements = []
                
                # 시트 데이터를 표로 변환
                table_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        table_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if table_data:
                    elements.append(DocumentElement(
                        ElementType.TABLE,
                        table_data,
                        {"table_type": "xlsx_sheet", "sheet_name": sheet_name}
                    ))
                
                processed_sheets.append(ProcessedPage(
                    sheet_num, elements, "sheet",
                    {"file_type": "xlsx", "sheet_name": sheet_name}
                ))
            
            wb.close()
            return processed_sheets
            
        except Exception as e:
            print(f"XLSX 처리 오류: {e}")
            return []
    
    def process_xls(self, file_path: str) -> List[ProcessedPage]:
        """XLS 파일을 시트별로 처리"""
        try:
            # 파일 내용 확인 (HTML 파일인지 체크)
            with open(file_path, 'rb') as f:
                first_bytes = f.read(100)
                if b'<html' in first_bytes.lower() or b'<!doctype' in first_bytes.lower():
                    print(f"⚠️ 파일이 HTML 형식입니다. 텍스트로 처리합니다: {file_path}")
                    return self.process_html(file_path)
            
            wb = xlrd.open_workbook(file_path)
            processed_sheets = []
            
            for sheet_num, sheet_name in enumerate(wb.sheet_names(), 1):
                sheet = wb.sheet_by_name(sheet_name)
                elements = []
                
                # 시트 데이터를 표로 변환
                table_data = []
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        row_data.append(str(cell_value) if cell_value else "")
                    
                    if any(cell for cell in row_data):
                        table_data.append(row_data)
                
                if table_data:
                    elements.append(DocumentElement(
                        ElementType.TABLE,
                        table_data,
                        {"table_type": "xls_sheet", "sheet_name": sheet_name}
                    ))
                
                processed_sheets.append(ProcessedPage(
                    sheet_num, elements, "sheet",
                    {"file_type": "xls", "sheet_name": sheet_name}
                ))
            
            return processed_sheets
            
        except Exception as e:
            print(f"XLS 처리 오류: {e}")
            # XLS 파일이 아닌 경우 텍스트로 처리 시도
            try:
                print(f"🔄 텍스트 파일로 처리 시도: {file_path}")
                return self._process_html_as_text(file_path)
            except Exception as e2:
                print(f"텍스트 처리도 실패: {e2}")
                return []
    
    def _refine_text_with_unstructured(self, text: str) -> str:
        """각 행의 텍스트를 Unstructured로 정제"""
        try:
            from unstructured.partition.text import partition_text
            
            # 텍스트를 Unstructured로 정제
            elements = partition_text(text=text)
            
            # 정제된 텍스트 추출
            refined_parts = []
            for element in elements:
                if hasattr(element, 'text') and element.text:
                    refined_parts.append(element.text.strip())
            
            if refined_parts:
                refined_text = ' '.join(refined_parts)
                if len(refined_text) > len(text) * 0.8:  # 원본의 80% 이상이면 사용
                    return refined_text
            
            # Unstructured 실패 시 원본 반환
            return text
            
        except Exception as e:
            print(f"   ⚠️ Unstructured 정제 실패: {e}")
            return text

    def _process_html_table_by_rows(self, file_path: str) -> List[ProcessedPage]:
        """HTML 테이블을 행별로 분할하여 처리 (각 티켓을 개별 청크로)"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            print(f"🔍 HTML 파일 내용 길이: {len(content)}자")
            print(f"🔍 HTML 파일 첫 200자: {content[:200]}")
            
            # BeautifulSoup으로 HTML 테이블 파싱
            from bs4 import BeautifulSoup
            
            print("🔧 BeautifulSoup으로 HTML 테이블 파싱 시도...")
            soup = BeautifulSoup(content, 'html.parser')
            
            # 테이블 찾기 (모든 테이블 중 가장 큰 것 선택)
            tables = soup.find_all('table')
            if not tables:
                print("⚠️ 테이블을 찾을 수 없음, 일반 텍스트로 처리")
                raise Exception("No table found")
            
            # 가장 많은 행을 가진 테이블 찾기 (실제 데이터 테이블)
            table_info = [(i, t, len(t.find_all('tr'))) for i, t in enumerate(tables)]
            table_info.sort(key=lambda x: x[2], reverse=True)  # 행 수로 내림차순 정렬
            
            print(f"🔍 테이블 분석:")
            for i, (table_idx, table, row_count) in enumerate(table_info):
                print(f"   테이블 {table_idx+1}: {row_count}행")
            
            # 가장 많은 행을 가진 테이블 선택
            table = table_info[0][1]
            selected_table_idx = table_info[0][0]
            selected_row_count = table_info[0][2]
            print(f"✅ 선택된 테이블: 테이블 {selected_table_idx+1} ({selected_row_count}행)")
            
            # 테이블 행들 추출
            rows = table.find_all('tr')
            print(f"🔍 발견된 테이블 행 수: {len(rows)}")
            
            if len(rows) < 2:
                print("⚠️ 데이터 행이 없음, 일반 텍스트로 처리")
                raise Exception("No data rows found")
            
            # 헤더 행과 데이터 행 분리
            header_row = rows[0]
            data_rows = rows[1:]
            
            # 헤더 추출 (더 강력한 방법)
            headers = []
            for cell in header_row.find_all(['td', 'th']):
                header_text = cell.get_text(strip=True)
                if header_text:
                    headers.append(header_text)
            
            # 헤더가 비어있으면 첫 번째 데이터 행을 헤더로 사용
            if not headers and data_rows:
                print("⚠️ 헤더가 비어있음, 첫 번째 데이터 행을 헤더로 사용")
                first_data_row = data_rows[0]
                for cell in first_data_row.find_all(['td', 'th']):
                    header_text = cell.get_text(strip=True)
                    if header_text:
                        headers.append(header_text)
                data_rows = data_rows[1:]  # 첫 번째 데이터 행 제거
            
            # 여전히 헤더가 없으면 기본 헤더 사용
            if not headers:
                print("⚠️ 헤더를 찾을 수 없음, 기본 헤더 사용")
                headers = ['Column1', 'Column2', 'Column3', 'Column4', 'Column5']
            
            print(f"🔍 테이블 헤더: {headers}")
            print(f"🔍 데이터 행 수: {len(data_rows)}")
            
            # 각 행의 셀 수 확인
            if data_rows:
                first_data_cells = data_rows[0].find_all(['td', 'th'])
                print(f"🔍 첫 번째 데이터 행의 셀 수: {len(first_data_cells)}")
                print(f"🔍 첫 번째 데이터 행 내용: {[cell.get_text(strip=True) for cell in first_data_cells]}")
            
            # 각 데이터 행을 개별 청크로 처리
            processed_elements = []
            
            for row_idx, row in enumerate(data_rows):
                # 행의 셀들 추출
                cells = row.find_all(['td', 'th'])
                if len(cells) == 0:
                    continue
                
                # 셀 텍스트 추출
                row_data = []
                for cell in cells:
                    cell_text = cell.get_text(strip=True)
                    if cell_text:
                        row_data.append(cell_text)
                
                if len(row_data) == 0:
                    continue
                
                # 헤더와 데이터를 매핑하여 구조화된 텍스트 생성
                ticket_text = ""
                for i, cell_text in enumerate(row_data):
                    if i < len(headers):
                        ticket_text += f"{headers[i]}: {cell_text} "
                    else:
                        ticket_text += f"{cell_text} "
                
                ticket_text = ticket_text.strip()
                
                if len(ticket_text) > 10:  # 의미있는 데이터만 처리
                    # 각 행에 Unstructured 적용하여 정제
                    refined_text = self._refine_text_with_unstructured(ticket_text)
                    
                    print(f"   📝 행 {row_idx + 1}: {refined_text[:100]}...")
                    
                    # 각 티켓을 개별 요소로 생성
                    element = DocumentElement(
                        ElementType.TEXT,
                        refined_text,
                        {
                            "source": "html_table_row",
                            "row_index": row_idx + 1,
                            "headers": headers,
                            "cell_count": len(row_data),
                            "unstructured_refined": True
                        }
                    )
                    processed_elements.append(element)
            
            if len(processed_elements) == 0:
                print("⚠️ 처리된 행이 없음, 일반 텍스트로 처리")
                raise Exception("No processed rows")
            
            print(f"✅ {len(processed_elements)}개 티켓 행 처리 완료")
            
            # 결과를 임시 파일로 저장
            import tempfile
            import json
            from datetime import datetime
            temp_dir = "temp_results"
            os.makedirs(temp_dir, exist_ok=True)
            
            result = {
                "file_path": file_path,
                "file_type": "html",
                "content_type": "text_based",
                "processed_pages": [{
                    "page_number": 1,
                    "page_type": "table",
                    "elements": [{
                        "element_type": "text",
                        "content": elem.content,
                        "metadata": elem.metadata
                    } for elem in processed_elements],
                    "metadata": {"file_type": "html", "table_rows": len(processed_elements)}
                }],
                "total_pages": 1,
                "processing_timestamp": datetime.now().isoformat()
            }
            
            temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='_result.json', delete=False, dir=temp_dir)
            json.dump(result, temp_file, ensure_ascii=False, indent=2)
            temp_file.close()
            print(f"결과가 임시 저장되었습니다: {temp_file.name}")
            
            return [ProcessedPage(
                page_number=1, 
                elements=processed_elements, 
                page_type="table",
                metadata={"file_type": "html", "table_rows": len(processed_elements)}
            )]
            
        except Exception as e:
            print(f"❌ HTML 테이블 행별 처리 오류: {e}")
            import traceback
            print(f"❌ 상세 오류 정보:")
            traceback.print_exc()
            return []

    def _process_html_as_text(self, file_path: str) -> List[ProcessedPage]:
        """HTML 파일을 텍스트로 처리 (Jira 테이블 최적화)"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            print(f"🔍 HTML 파일 내용 길이: {len(content)}자")
            print(f"🔍 HTML 파일 첫 200자: {content[:200]}")
            
            # Unstructured 라이브러리로 정교한 텍스트 추출 시도
            try:
                from unstructured.partition.html import partition_html
                from unstructured.staging.base import elements_to_json
                
                print("🔧 Unstructured 라이브러리로 HTML 파싱 시도...")
                
                # HTML을 구조화된 요소로 분할
                elements = partition_html(text=content)
                
                # 텍스트만 추출
                text_parts = []
                for element in elements:
                    if hasattr(element, 'text') and element.text:
                        text_parts.append(element.text.strip())
                
                unstructured_text = ' '.join(text_parts)
                
                print(f"🔍 Unstructured 추출 결과: {len(unstructured_text)}자")
                print(f"🔍 Unstructured 첫 200자: {unstructured_text[:200]}")
                
                # Unstructured 결과가 충분한지 확인 (최소 500자 이상)
                if unstructured_text and len(unstructured_text) > 500:
                    # Unstructured 결과가 좋으면 사용
                    text_content = unstructured_text
                    print("✅ Unstructured 라이브러리 사용")
                else:
                    print(f"⚠️ Unstructured 결과 부족 ({len(unstructured_text)}자), BeautifulSoup으로 폴백")
                    raise Exception("Unstructured 결과가 부족함")
                    
            except Exception as e:
                print(f"⚠️ Unstructured 라이브러리 실패, BeautifulSoup으로 폴백: {e}")
                
                # BeautifulSoup으로 테이블 행별 분할 시도
                try:
                    from bs4 import BeautifulSoup
                    
                    print("🔧 BeautifulSoup으로 HTML 테이블 파싱 시도...")
                    soup = BeautifulSoup(content, 'html.parser')
                    
                    # 테이블 찾기
                    table = soup.find('table')
                    if table:
                        # 테이블 행들 추출
                        rows = table.find_all('tr')
                        print(f"🔍 발견된 테이블 행 수: {len(rows)}")
                        
                        if len(rows) >= 2:
                            # 헤더 행과 데이터 행 분리
                            header_row = rows[0]
                            data_rows = rows[1:]
                            
                            # 헤더 추출
                            headers = []
                            for cell in header_row.find_all(['td', 'th']):
                                header_text = cell.get_text(strip=True)
                                if header_text:
                                    headers.append(header_text)
                            
                            print(f"🔍 테이블 헤더: {headers}")
                            print(f"🔍 데이터 행 수: {len(data_rows)}")
                            
                            # 각 데이터 행을 개별 청크로 처리
                            processed_elements = []
                            
                            for row_idx, row in enumerate(data_rows):
                                # 행의 셀들 추출
                                cells = row.find_all(['td', 'th'])
                                if len(cells) == 0:
                                    continue
                                
                                # 셀 텍스트 추출
                                row_data = []
                                for cell in cells:
                                    cell_text = cell.get_text(strip=True)
                                    if cell_text:
                                        row_data.append(cell_text)
                                
                                if len(row_data) == 0:
                                    continue
                                
                                # 헤더와 데이터를 매핑하여 구조화된 텍스트 생성
                                ticket_text = ""
                                for i, cell_text in enumerate(row_data):
                                    if i < len(headers):
                                        ticket_text += f"{headers[i]}: {cell_text} "
                                    else:
                                        ticket_text += f"{cell_text} "
                                
                                ticket_text = ticket_text.strip()
                                
                                if len(ticket_text) > 10:  # 의미있는 데이터만 처리
                                    print(f"   📝 행 {row_idx + 1}: {ticket_text[:100]}...")
                                    
                                    # 각 티켓을 개별 요소로 생성
                                    element = DocumentElement(
                                        ElementType.TEXT,
                                        ticket_text,
                                        {
                                            "source": "html_table_row",
                                            "row_index": row_idx + 1,
                                            "headers": headers,
                                            "cell_count": len(row_data)
                                        }
                                    )
                                    processed_elements.append(element)
                            
                            if len(processed_elements) > 0:
                                print(f"✅ {len(processed_elements)}개 티켓 행 처리 완료")
                                
                                # 결과를 임시 파일로 저장
                                import tempfile
                                import json
                                from datetime import datetime
                                temp_dir = "temp_results"
                                os.makedirs(temp_dir, exist_ok=True)
                                
                                result = {
                                    "file_path": file_path,
                                    "file_type": "html",
                                    "content_type": "text_based",
                                    "processed_pages": [{
                                        "page_number": 1,
                                        "page_type": "table",
                                        "elements": [{
                                            "element_type": "text",
                                            "content": elem.content,
                                            "metadata": elem.metadata
                                        } for elem in processed_elements],
                                        "metadata": {"file_type": "html", "table_rows": len(processed_elements)}
                                    }],
                                    "total_pages": 1,
                                    "processing_timestamp": datetime.now().isoformat()
                                }
                                
                                temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='_result.json', delete=False, dir=temp_dir)
                                json.dump(result, temp_file, ensure_ascii=False, indent=2)
                                temp_file.close()
                                print(f"결과가 임시 저장되었습니다: {temp_file.name}")
                                
                                return [ProcessedPage(
                                    1, 
                                    processed_elements, 
                                    "table",
                                    {"file_type": "html", "table_rows": len(processed_elements)}
                                )]
                            else:
                                print("⚠️ 처리된 행이 없음, 정규식으로 처리")
                                raise Exception("No processed rows")
                        else:
                            print("⚠️ 테이블 데이터 행이 없음, 정규식으로 처리")
                            raise Exception("No data rows")
                    else:
                        print("⚠️ 테이블을 찾을 수 없음, 정규식으로 처리")
                        raise Exception("No table found")
                        
                except Exception as e2:
                    print(f"⚠️ BeautifulSoup 처리도 실패, 정규식으로 처리: {e2}")
                
                # HTML 태그 제거 (더 강력한 방법)
                import re
                
                # 1. CSS 스타일 블록 제거 (Jira 테이블 스타일)
                text_content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL | re.IGNORECASE)
                
                # 2. HTML 태그 제거
                text_content = re.sub(r'<[^>]+>', '', text_content)
                
                # 3. HTML 엔티티 디코딩
                text_content = text_content.replace('&nbsp;', ' ')
                text_content = text_content.replace('&amp;', '&')
                text_content = text_content.replace('&lt;', '<')
                text_content = text_content.replace('&gt;', '>')
                text_content = text_content.replace('&quot;', '"')
                text_content = text_content.replace('&#39;', "'")
                
                # 4. Jira 테이블 관련 CSS 속성 제거
                text_content = re.sub(r'mso-[^;]+;?', '', text_content)
                text_content = re.sub(r'@page[^}]+}', '', text_content)
                text_content = re.sub(r'body\s*{[^}]*}', '', text_content)
                text_content = re.sub(r'table\s*{[^}]*}', '', text_content)
                
                # 5. 연속 공백 및 특수문자 정리
                text_content = re.sub(r'\s+', ' ', text_content)
                text_content = re.sub(r'[{}();]', ' ', text_content)
                text_content = re.sub(r'\s+', ' ', text_content).strip()
                
                # 6. 의미있는 텍스트만 추출 (너무 짧은 단어 제거)
                words = text_content.split()
                meaningful_words = [word for word in words if len(word) > 2 or word.isdigit()]
                text_content = ' '.join(meaningful_words)
            
            print(f"🔍 최종 텍스트 길이: {len(text_content)}자")
            print(f"🔍 최종 텍스트 첫 200자: {text_content[:200]}")
            
            if text_content and len(text_content) > 50:  # 최소 길이 체크 (50자 이상)
                elements = [DocumentElement(
                    ElementType.TEXT,
                    text_content,
                    {"content_type": "jira_table_text", "original_length": len(content), "processed": True}
                )]
                
                return [ProcessedPage(
                    1, elements, "text",
                    {"file_type": "jira_table", "original_type": "xls", "text_length": len(text_content)}
                )]
            else:
                print(f"⚠️ 텍스트가 너무 짧거나 비어있음: {len(text_content)}자")
                return []
                
        except Exception as e:
            print(f"HTML 텍스트 처리 오류: {e}")
            return []
    
    def process_html(self, file_path: str) -> List[ProcessedPage]:
        """HTML 파일을 처리 (테이블 행별 분할 우선 시도)"""
        # 먼저 테이블 행별 분할 시도
        try:
            result = self._process_html_table_by_rows(file_path)
            if result and len(result) > 0 and len(result[0].elements) > 0:
                print("✅ HTML 테이블 행별 분할 성공")
                return result
        except Exception as e:
            print(f"⚠️ HTML 테이블 행별 분할 실패, 일반 텍스트로 처리: {e}")
        
        # 폴백: 일반 HTML 텍스트 처리
        return self._process_html_as_text(file_path)
    
    def _extract_pptx_image(self, shape, slide_num: int, element_index: int) -> str:
        """PPTX에서 이미지 추출"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # 임시 디렉토리에 저장
            temp_dir = "temp_images"
            os.makedirs(temp_dir, exist_ok=True)
            
            image_filename = f"pptx_slide{slide_num}_img{element_index}.png"
            image_path = os.path.join(temp_dir, image_filename)
            
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            
            return image_path
            
        except Exception as e:
            raise Exception(f"이미지 추출 실패: {e}")
    
    def _process_xml(self, file_path: str) -> List[ProcessedPage]:
        """
        XML 파일을 처리하여 ProcessedPage 리스트를 반환
        
        Args:
            file_path: XML 파일 경로
            
        Returns:
            XML 파싱 결과를 포함한 ProcessedPage 리스트
        """
        try:
            print(f"🔍 XML 파일 처리 시작: {file_path}")
            
            # XML 파싱
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            print(f"📋 XML 루트 요소: {root.tag}")
            
            processed_elements = []
            
            # RSS 형식인지 확인
            if root.tag == 'rss':
                processed_elements = self._process_rss_xml(root, file_path)
            else:
                # 일반 XML 처리
                processed_elements = self._process_generic_xml(root, file_path)
            
            if not processed_elements:
                print("⚠️ 처리된 요소가 없음, 전체 XML을 텍스트로 처리")
                processed_elements = self._process_xml_as_text(file_path)
            
            print(f"✅ XML 처리 완료: {len(processed_elements)}개 요소")
            
            # 결과를 ProcessedPage로 변환
            page = ProcessedPage(
                page_number=1,
                page_type="xml_content",
                elements=processed_elements,
                metadata={
                    "file_type": "xml",
                    "root_element": root.tag,
                    "total_elements": len(processed_elements),
                    "processing_method": "xml_parser"
                }
            )
            
            return [page]
            
        except ET.ParseError as e:
            print(f"❌ XML 파싱 오류: {e}")
            # 파싱 실패 시 텍스트로 처리
            return self._process_xml_as_text_fallback(file_path)
        except Exception as e:
            print(f"❌ XML 처리 오류: {e}")
            return self._process_xml_as_text_fallback(file_path)
    
    def _process_rss_xml(self, root, file_path: str) -> List[DocumentElement]:
        """RSS XML 처리 (JIRA RSS 특화)"""
        processed_elements = []
        
        # 채널 정보 추출
        channel = root.find('channel')
        if channel is not None:
            title_elem = channel.find('title')
            channel_title = title_elem.text if title_elem is not None else "Unknown"
            
            print(f"📋 RSS 채널: {channel_title}")
            
            # 이슈 정보 추출
            issue_elem = channel.find('issue')
            if issue_elem is not None:
                total_issues = issue_elem.get('total', '0')
                print(f"📋 총 이슈 수: {total_issues}")
            
            # 각 아이템(이슈) 처리
            items = channel.findall('item')
            print(f"📋 처리할 아이템 수: {len(items)}")
            
            for idx, item in enumerate(items):
                try:
                    # 이슈 정보 추출
                    issue_data = self._extract_jira_issue_data(item)
                    
                    if issue_data:
                        # 구조화된 텍스트 생성
                        structured_text = self._create_structured_issue_text(issue_data)
                        
                        # Unstructured로 정제
                        refined_text = self._refine_text_with_unstructured(structured_text)
                        
                        # DocumentElement 생성
                        element = DocumentElement(
                            ElementType.TEXT,
                            refined_text,
                            {
                                "source": "jira_rss_item",
                                "issue_key": issue_data.get('key', ''),
                                "issue_type": issue_data.get('type', ''),
                                "status": issue_data.get('status', ''),
                                "priority": issue_data.get('priority', ''),
                                "assignee": issue_data.get('assignee', ''),
                                "reporter": issue_data.get('reporter', ''),
                                "created": issue_data.get('created', ''),
                                "updated": issue_data.get('updated', ''),
                                "item_index": idx + 1,
                                "unstructured_refined": True
                            }
                        )
                        processed_elements.append(element)
                        
                        if (idx + 1) % 100 == 0:
                            print(f"   📝 {idx + 1}개 이슈 처리 완료...")
                
                except Exception as e:
                    print(f"⚠️ 아이템 {idx + 1} 처리 오류: {e}")
                    continue
        
        return processed_elements
    
    def _extract_jira_issue_data(self, item) -> Dict[str, str]:
        """JIRA 이슈 데이터 추출"""
        issue_data = {}
        
        try:
            # 기본 정보
            title_elem = item.find('title')
            if title_elem is not None:
                issue_data['title'] = title_elem.text or ''
            
            key_elem = item.find('key')
            if key_elem is not None:
                issue_data['key'] = key_elem.text or ''
            
            summary_elem = item.find('summary')
            if summary_elem is not None:
                issue_data['summary'] = summary_elem.text or ''
            
            # 설명 (HTML 디코딩)
            description_elem = item.find('description')
            if description_elem is not None and description_elem.text:
                # HTML 엔티티 디코딩
                description = html.unescape(description_elem.text)
                # HTML 태그 제거 (간단한 방법)
                import re
                description = re.sub(r'<[^>]+>', '', description)
                issue_data['description'] = description.strip()
            
            # 프로젝트 정보
            project_elem = item.find('project')
            if project_elem is not None:
                issue_data['project'] = project_elem.text or ''
                issue_data['project_key'] = project_elem.get('key', '')
            
            # 이슈 타입
            type_elem = item.find('type')
            if type_elem is not None:
                issue_data['type'] = type_elem.text or ''
            
            # 상태
            status_elem = item.find('status')
            if status_elem is not None:
                issue_data['status'] = status_elem.text or ''
            
            # 우선순위
            priority_elem = item.find('priority')
            if priority_elem is not None:
                issue_data['priority'] = priority_elem.text or ''
            
            # 담당자
            assignee_elem = item.find('assignee')
            if assignee_elem is not None:
                issue_data['assignee'] = assignee_elem.text or ''
            
            # 보고자
            reporter_elem = item.find('reporter')
            if reporter_elem is not None:
                issue_data['reporter'] = reporter_elem.text or ''
            
            # 날짜 정보
            created_elem = item.find('created')
            if created_elem is not None:
                issue_data['created'] = created_elem.text or ''
            
            updated_elem = item.find('updated')
            if updated_elem is not None:
                issue_data['updated'] = updated_elem.text or ''
            
            resolved_elem = item.find('resolved')
            if resolved_elem is not None:
                issue_data['resolved'] = resolved_elem.text or ''
            
            # 링크
            link_elem = item.find('link')
            if link_elem is not None:
                issue_data['link'] = link_elem.text or ''
            
        except Exception as e:
            print(f"⚠️ 이슈 데이터 추출 오류: {e}")
        
        return issue_data
    
    def _create_structured_issue_text(self, issue_data: Dict[str, str]) -> str:
        """구조화된 이슈 텍스트 생성"""
        text_parts = []
        
        # 이슈 키와 제목
        if issue_data.get('key'):
            text_parts.append(f"이슈 키: {issue_data['key']}")
        if issue_data.get('title'):
            text_parts.append(f"제목: {issue_data['title']}")
        
        # 요약
        if issue_data.get('summary'):
            text_parts.append(f"요약: {issue_data['summary']}")
        
        # 설명
        if issue_data.get('description'):
            text_parts.append(f"설명: {issue_data['description']}")
        
        # 메타데이터
        metadata_parts = []
        if issue_data.get('type'):
            metadata_parts.append(f"타입: {issue_data['type']}")
        if issue_data.get('status'):
            metadata_parts.append(f"상태: {issue_data['status']}")
        if issue_data.get('priority'):
            metadata_parts.append(f"우선순위: {issue_data['priority']}")
        if issue_data.get('assignee'):
            metadata_parts.append(f"담당자: {issue_data['assignee']}")
        if issue_data.get('reporter'):
            metadata_parts.append(f"보고자: {issue_data['reporter']}")
        if issue_data.get('created'):
            metadata_parts.append(f"생성일: {issue_data['created']}")
        if issue_data.get('updated'):
            metadata_parts.append(f"수정일: {issue_data['updated']}")
        
        if metadata_parts:
            text_parts.append("메타데이터: " + ", ".join(metadata_parts))
        
        return "\n".join(text_parts)
    
    def _process_generic_xml(self, root, file_path: str) -> List[DocumentElement]:
        """일반 XML 처리"""
        processed_elements = []
        
        def extract_text_recursive(element, depth: int = 0) -> str:
            """XML 요소에서 텍스트 재귀적으로 추출"""
            text_parts = []
            
            # 현재 요소의 텍스트
            if element.text and element.text.strip():
                text_parts.append(element.text.strip())
            
            # 자식 요소들 처리
            for child in element:
                child_text = extract_text_recursive(child, depth + 1)
                if child_text:
                    text_parts.append(child_text)
            
            # 현재 요소의 tail 텍스트
            if element.tail and element.tail.strip():
                text_parts.append(element.tail.strip())
            
            return " ".join(text_parts)
        
        # 루트 요소에서 텍스트 추출
        full_text = extract_text_recursive(root)
        
        if full_text.strip():
            # Unstructured로 정제
            refined_text = self._refine_text_with_unstructured(full_text)
            
            element = DocumentElement(
                ElementType.TEXT,
                refined_text,
                {
                    "source": "generic_xml",
                    "root_element": root.tag,
                    "unstructured_refined": True
                }
            )
            processed_elements.append(element)
        
        return processed_elements
    
    def _process_xml_as_text(self, file_path: str) -> List[DocumentElement]:
        """XML을 일반 텍스트로 처리"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # HTML 엔티티 디코딩
            content = html.unescape(content)
            
            # Unstructured로 정제
            refined_text = self._refine_text_with_unstructured(content)
            
            element = DocumentElement(
                ElementType.TEXT,
                refined_text,
                {
                    "source": "xml_as_text",
                    "unstructured_refined": True
                }
            )
            return [element]
            
        except Exception as e:
            print(f"❌ XML 텍스트 처리 오류: {e}")
            return []
    
    def _process_xml_as_text_fallback(self, file_path: str) -> List[ProcessedPage]:
        """XML 파싱 실패 시 텍스트로 처리하는 fallback"""
        try:
            elements = self._process_xml_as_text(file_path)
            
            page = ProcessedPage(
                page_number=1,
                page_type="xml_text_fallback",
                elements=elements,
                metadata={
                    "file_type": "xml",
                    "processing_method": "text_fallback",
                    "error": "xml_parse_failed"
                }
            )
            
            return [page]
            
        except Exception as e:
            print(f"❌ XML fallback 처리 오류: {e}")
            return []

class LayoutBasedProcessor:
    """Layout-based 파일 처리기 (PDF 변환 후 이미지 처리)"""
    
    def __init__(self, azure_processor: AzureOpenAIImageProcessor):
        self.azure_processor = azure_processor
    
    def process_file(self, file_path: str, doc_type: DocumentType) -> List[ProcessedPage]:
        """Layout-based 파일을 PDF로 변환 후 처리"""
        try:
            # PDF로 변환
            pdf_path = self._convert_to_pdf(file_path, doc_type)
            
            # PDF를 이미지로 변환하여 처리
            return self._process_pdf_as_images(pdf_path)
            
        except Exception as e:
            print(f"Layout-based 처리 오류: {e}")
            return []
    
    def _convert_to_pdf(self, file_path: str, doc_type: DocumentType) -> str:
        """파일을 PDF로 변환"""
        try:
            if doc_type == DocumentType.PDF:
                return file_path
            
            # 임시 PDF 파일 경로
            temp_dir = "temp_pdfs"
            os.makedirs(temp_dir, exist_ok=True)
            
            pdf_filename = f"converted_{Path(file_path).stem}.pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)
            
            if doc_type == DocumentType.DOCX:
                self._convert_docx_to_pdf(file_path, pdf_path)
            elif doc_type == DocumentType.PPTX:
                self._convert_pptx_to_pdf(file_path, pdf_path)
            elif doc_type == DocumentType.XLSX:
                self._convert_xlsx_to_pdf(file_path, pdf_path)
            
            return pdf_path
            
        except Exception as e:
            raise Exception(f"PDF 변환 실패: {e}")
    
    def _convert_docx_to_pdf(self, docx_path: str, pdf_path: str):
        """DOCX를 PDF로 변환 (LibreOffice 사용)"""
        try:
            import subprocess
            # LibreOffice를 사용하여 변환
            cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", 
                   os.path.dirname(pdf_path), docx_path]
            subprocess.run(cmd, check=True, capture_output=True)
            
            # 변환된 파일 이름 변경
            converted_path = docx_path.replace('.docx', '.pdf')
            if os.path.exists(converted_path):
                shutil.move(converted_path, pdf_path)
                
        except Exception as e:
            print(f"LibreOffice 변환 실패, 대안 방법 사용: {e}")
            # 대안: python-docx2pdf 사용
            try:
                from docx2pdf import convert
                convert(docx_path, pdf_path)
            except ImportError:
                raise Exception("PDF 변환을 위한 라이브러리가 설치되지 않았습니다.")
    
    def _convert_pptx_to_pdf(self, pptx_path: str, pdf_path: str):
        """PPTX를 PDF로 변환"""
        try:
            from pptx2pdf import convert
            convert(pptx_path, pdf_path)
        except ImportError:
            # 대안: LibreOffice 사용
            self._convert_docx_to_pdf(pptx_path, pdf_path)
    
    def _convert_xlsx_to_pdf(self, xlsx_path: str, pdf_path: str):
        """XLSX를 PDF로 변환"""
        try:
            # pandas를 사용하여 HTML로 변환 후 PDF 변환
            df = pd.read_excel(xlsx_path)
            html_path = xlsx_path.replace('.xlsx', '.html')
            df.to_html(html_path)
            
            # HTML을 PDF로 변환 (weasyprint 사용)
            try:
                from weasyprint import HTML
                HTML(html_path).write_pdf(pdf_path)
                os.remove(html_path)
            except ImportError:
                # 대안: LibreOffice 사용
                self._convert_docx_to_pdf(xlsx_path, pdf_path)
                
        except Exception as e:
            print(f"XLSX 변환 실패, LibreOffice 사용: {e}")
            self._convert_docx_to_pdf(xlsx_path, pdf_path)
    
    def _process_pdf_as_images(self, pdf_path: str) -> List[ProcessedPage]:
        """PDF를 이미지로 변환하여 GPT Vision으로 처리"""
        try:
            doc = fitz.open(pdf_path)
            processed_pages = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 페이지를 이미지로 변환
                mat = fitz.Matrix(2.0, 2.0)  # 2배 확대
                pix = page.get_pixmap(matrix=mat)
                
                # 이미지 저장
                temp_dir = "temp_images"
                os.makedirs(temp_dir, exist_ok=True)
                
                image_filename = f"pdf_page{page_num + 1}.png"
                image_path = os.path.join(temp_dir, image_filename)
                pix.save(image_path)
                
                # GPT Vision으로 이미지 처리
                prompt = "이 페이지의 모든 텍스트 내용을 추출하고, 표나 이미지가 있다면 설명해주세요."
                text_content = self.azure_processor.image_to_text(image_path, prompt)
                
                # 결과를 텍스트 요소로 저장
                elements = [DocumentElement(
                    ElementType.TEXT,
                    text_content,
                    {"page_number": page_num + 1, "source": "gpt_vision"}
                )]
                
                processed_pages.append(ProcessedPage(
                    page_num + 1, elements, "page",
                    {"file_type": "pdf", "processing_method": "gpt_vision"}
                ))
            
            doc.close()
            return processed_pages
            
        except Exception as e:
            raise Exception(f"PDF 이미지 처리 실패: {e}")

class FileProcessor:
    """메인 파일 처리기"""
    
    def __init__(self, azure_processor: AzureOpenAIImageProcessor):
        self.azure_processor = azure_processor
        self.text_processor = TextBasedProcessor(azure_processor)
        self.layout_processor = LayoutBasedProcessor(azure_processor)
        self.temp_storage = {}
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """파일을 처리하고 결과를 반환"""
        try:
            # 1. 파일 타입 판별
            doc_type = FileTypeDetector.detect_file_type(file_path)
            content_type = FileTypeDetector.detect_content_type(file_path, doc_type)
            
            print(f"파일 타입: {doc_type.value}, 콘텐츠 타입: {content_type.value}")
            
            # 2. 콘텐츠 타입에 따른 처리
            if content_type == ContentType.TEXT_BASED:
                processed_pages = self._process_text_based(file_path, doc_type)
            else:
                processed_pages = self._process_layout_based(file_path, doc_type)
            
            # 3. 결과를 메타데이터와 함께 임시 저장
            result = {
                "file_path": file_path,
                "file_type": doc_type.value,
                "content_type": content_type.value,
                "processed_pages": [page.to_dict() for page in processed_pages],
                "total_pages": len(processed_pages),
                "processing_timestamp": str(pd.Timestamp.now())
            }
            
            # 임시 저장
            self._save_to_temp_storage(file_path, result)
            
            return result

        except Exception as e:
            print(f"파일 처리 오류: {e}")
            return {"error": str(e)}
    
    def _process_text_based(self, file_path: str, doc_type: DocumentType) -> List[ProcessedPage]:
        """Text-based 파일 처리"""
        if doc_type == DocumentType.DOCX:
            return self.text_processor.process_docx(file_path)
        elif doc_type == DocumentType.PPTX:
            return self.text_processor.process_pptx(file_path)
        elif doc_type == DocumentType.XLSX:
            return self.text_processor.process_xlsx(file_path)
        elif doc_type == DocumentType.XLS:
            return self.text_processor.process_xls(file_path)
        elif doc_type == DocumentType.HTML:
            return self.text_processor.process_html(file_path)
        elif doc_type == DocumentType.XML:
            return self.text_processor._process_xml(file_path)
        elif doc_type == DocumentType.PDF:
            # PDF는 별도 처리
            return self._process_pdf_text_based(file_path)
        else:
            return []

    def _process_layout_based(self, file_path: str, doc_type: DocumentType) -> List[ProcessedPage]:
        """Layout-based 파일 처리"""
        return self.layout_processor.process_file(file_path, doc_type)
    
    def _process_pdf_text_based(self, file_path: str) -> List[ProcessedPage]:
        """Text-based PDF 처리"""
        try:
            doc = fitz.open(file_path)
            processed_pages = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                elements = []
                
                # 텍스트 블록 추출
                text_blocks = page.get_text("blocks")
                for block in text_blocks:
                    x0, y0, x1, y1, text, block_no, block_type = block[:7]
                    
                    if block_type == 0 and text.strip():  # 텍스트 블록
                            elements.append(DocumentElement(
                                ElementType.TEXT, 
                                text.strip(),
                            {"page": page_num + 1, "bbox": (x0, y0, x1, y1)}
                        ))
                
                if elements:
                    processed_pages.append(ProcessedPage(
                        page_num + 1, elements, "page",
                        {"file_type": "pdf", "processing_method": "text_extraction"}
                        ))
            
            doc.close()
            return processed_pages
            
        except Exception as e:
            print(f"PDF 텍스트 처리 오류: {e}")
            return []
    
    def _save_to_temp_storage(self, file_path: str, result: Dict[str, Any]):
        """결과를 임시 저장소에 저장"""
        file_key = Path(file_path).name
        self.temp_storage[file_key] = result
        
        # JSON 파일로도 저장
        temp_dir = "temp_results"
        os.makedirs(temp_dir, exist_ok=True)
        
        json_filename = f"{Path(file_path).stem}_result.json"
        json_path = os.path.join(temp_dir, json_filename)
        
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        
        print(f"결과가 임시 저장되었습니다: {json_path}")
    
    def get_temp_result(self, file_name: str) -> Optional[Dict[str, Any]]:
        """임시 저장된 결과 조회"""
        return self.temp_storage.get(file_name)
    
    def clear_temp_storage(self):
        """임시 저장소 정리"""
        self.temp_storage.clear()
    
    def process_with_structured_chunking(self, file_path: str) -> List[StructuredChunk]:
        """
        구조적 청킹을 사용하여 파일 처리
        
        Args:
            file_path: 처리할 파일 경로
            
        Returns:
            구조화된 청크 리스트
        """
        try:
            file_extension = Path(file_path).suffix.lower()
            file_name = Path(file_path).name
            
            print(f"🔧 구조적 청킹 시작: {file_name}")
            
            # Jira HTML 파일인 경우
            if file_extension in ['.xls', '.html'] and self._is_jira_html(file_path):
                print("📋 Jira HTML 파일 감지 - 구조적 청킹 적용")
                with open(file_path, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                chunker = JiraStructuredChunker()
                return chunker.chunk_jira_html(html_content, file_name)
            
            # CSV 파일인 경우
            elif file_extension == '.csv':
                print("📋 CSV 파일 감지 - 구조적 청킹 적용")
                try:
                    df = pd.read_csv(file_path)
                    csv_data = df.to_dict('records')
                    
                    chunker = JiraStructuredChunker()
                    return chunker.chunk_csv_data(csv_data, file_name)
                except Exception as e:
                    print(f"❌ CSV 구조적 청킹 실패: {str(e)}")
                    return []
            
            else:
                print(f"⚠️ 구조적 청킹 미지원 파일: {file_extension}")
                return []
                
        except Exception as e:
            print(f"❌ 구조적 청킹 실패: {str(e)}")
            return []
    
    def _is_jira_html(self, file_path: str) -> bool:
        """파일이 Jira HTML인지 확인"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read(1000)  # 첫 1000자만 읽기
            
            # Jira HTML의 특징적인 패턴들 확인
            jira_indicators = [
                'Jira',
                'issuerow',
                'issue-link',
                'BTVO-',
                '프로젝트:',
                '키:',
                '요약:'
            ]
            
            return any(indicator in content for indicator in jira_indicators)
            
        except Exception:
            return False
        
        # 임시 파일들도 정리
        temp_dirs = ["temp_images", "temp_pdfs", "temp_results"]
        for temp_dir in temp_dirs:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                print(f"임시 디렉토리 정리됨: {temp_dir}")
