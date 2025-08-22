"""페이지별 콘텐츠 분석기"""

import os
import tempfile
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook
from PIL import Image
import io

from module.types import DocumentType, ContentType, ElementType, ContentMetrics, ClassificationResult


class PageContentAnalyzer:
    """개별 페이지/슬라이드의 콘텐츠 타입을 분석하는 클래스"""
    
    def __init__(self):
        # 페이지별 분석을 위한 임계값 (파일 전체 분석보다 더 엄격)
        self.thresholds = {
            "min_text_length": 50,  # 최소 텍스트 길이
            "max_image_ratio": 0.3,  # 최대 이미지 비율
            "min_text_density": 0.15,  # 최소 텍스트 밀도
            "layout_complexity_threshold": 0.4,  # 레이아웃 복잡도 임계값
        }
    
    def analyze_pdf_page(self, pdf_path: str, page_num: int) -> ClassificationResult:
        """PDF 특정 페이지 분석"""
        try:
            doc = fitz.open(pdf_path)
            page = doc[page_num - 1]  # 0-based index
            
            # 텍스트 추출 및 분석
            text = page.get_text()
            text_blocks = page.get_text("dict")
            
            # 이미지 정보 추출
            image_list = page.get_images()
            
            # 메트릭 계산
            metrics = ContentMetrics()
            metrics.text_length = len(text.strip())
            metrics.word_count = len(text.split())
            metrics.image_count = len(image_list)
            metrics.has_selectable_text = bool(text.strip())
            
            # 텍스트 밀도 계산 (페이지 크기 대비)
            page_area = page.rect.width * page.rect.height
            text_area = sum(block.get("bbox", [0, 0, 0, 0])[2] * block.get("bbox", [0, 0, 0, 0])[3] 
                          for block in text_blocks.get("blocks", []) if "lines" in block)
            metrics.text_density = text_area / page_area if page_area > 0 else 0
            
            # 레이아웃 복잡도 분석
            layout_complexity = self._analyze_pdf_layout_complexity(text_blocks)
            
            doc.close()
            
            # 분류 결정
            content_type, confidence, reasoning = self._classify_page_content(
                metrics, layout_complexity
            )
            
            return ClassificationResult(
                content_type=content_type,
                confidence=confidence,
                reasoning=reasoning,
                metrics=metrics,
                feature_scores={"layout_complexity": layout_complexity}
            )
            
        except Exception as e:
            print(f"PDF 페이지 {page_num} 분석 오류: {e}")
            return ClassificationResult(
                content_type=ContentType.LAYOUT_BASED,
                confidence=0.5,
                reasoning=[f"분석 오류로 인한 기본값: {str(e)}"],
                metrics=ContentMetrics()
            )
    
    def analyze_docx_page(self, docx_path: str, page_num: int) -> ClassificationResult:
        """DOCX 가상 페이지 분석 (문단 기준으로 나눔)"""
        try:
            doc = Document(docx_path)
            paragraphs = doc.paragraphs
            
            # 페이지 당 약 10-15개 문단으로 가정하여 분할
            paragraphs_per_page = 12
            start_idx = (page_num - 1) * paragraphs_per_page
            end_idx = min(start_idx + paragraphs_per_page, len(paragraphs))
            
            page_paragraphs = paragraphs[start_idx:end_idx]
            
            # 텍스트 분석
            text_content = "\n".join([p.text for p in page_paragraphs])
            
            # 이미지 분석 (해당 문단 범위에서)
            image_count = 0
            for para in page_paragraphs:
                for run in para.runs:
                    if run._element.xpath('.//pic:pic'):
                        image_count += 1
            
            # 메트릭 계산
            metrics = ContentMetrics()
            metrics.text_length = len(text_content.strip())
            metrics.word_count = len(text_content.split())
            metrics.paragraph_count = len([p for p in page_paragraphs if p.text.strip()])
            metrics.image_count = image_count
            
            # 스타일 분석
            style_diversity = len(set(p.style.name for p in page_paragraphs if p.style))
            metrics.font_variation = style_diversity  # font_diversity 대신 font_variation 사용
            
            # 분류 결정
            content_type, confidence, reasoning = self._classify_page_content(metrics)
            
            return ClassificationResult(
                content_type=content_type,
                confidence=confidence,
                reasoning=reasoning,
                metrics=metrics
            )
            
        except Exception as e:
            print(f"DOCX 페이지 {page_num} 분석 오류: {e}")
            return ClassificationResult(
                content_type=ContentType.LAYOUT_BASED,
                confidence=0.5,
                reasoning=[f"분석 오류로 인한 기본값: {str(e)}"],
                metrics=ContentMetrics()
            )
    
    def analyze_pptx_slide(self, pptx_path: str, slide_num: int) -> ClassificationResult:
        """PPTX 특정 슬라이드 분석"""
        try:
            prs = Presentation(pptx_path)
            if slide_num > len(prs.slides):
                raise IndexError(f"슬라이드 {slide_num}이 존재하지 않습니다")
            
            slide = prs.slides[slide_num - 1]  # 0-based index
            
            # 텍스트 및 도형 분석
            text_content = ""
            image_count = 0
            text_box_count = 0
            shape_count = 0
            
            for shape in slide.shapes:
                shape_count += 1
                
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text + "\n"
                    text_box_count += 1
                
                if shape.shape_type == 13:  # 이미지
                    image_count += 1
            
            # 메트릭 계산
            metrics = ContentMetrics()
            metrics.text_length = len(text_content.strip())
            metrics.word_count = len(text_content.split())
            metrics.image_count = image_count
            metrics.shape_count = shape_count
            metrics.text_box_count = text_box_count
            
            # 레이아웃 복잡도 (도형 수 기준)
            layout_complexity = min(shape_count / 10.0, 1.0)  # 10개 이상이면 복잡
            
            # 분류 결정
            content_type, confidence, reasoning = self._classify_page_content(
                metrics, layout_complexity
            )
            
            return ClassificationResult(
                content_type=content_type,
                confidence=confidence,
                reasoning=reasoning,
                metrics=metrics,
                feature_scores={"layout_complexity": layout_complexity}
            )
            
        except Exception as e:
            print(f"PPTX 슬라이드 {slide_num} 분석 오류: {e}")
            return ClassificationResult(
                content_type=ContentType.LAYOUT_BASED,
                confidence=0.5,
                reasoning=[f"분석 오류로 인한 기본값: {str(e)}"],
                metrics=ContentMetrics()
            )
    
    def analyze_xlsx_sheet(self, xlsx_path: str, sheet_num: int) -> ClassificationResult:
        """XLSX 특정 시트 분석"""
        try:
            wb = load_workbook(xlsx_path, data_only=True)
            sheet_names = wb.sheetnames
            
            if sheet_num > len(sheet_names):
                raise IndexError(f"시트 {sheet_num}이 존재하지 않습니다")
            
            ws = wb[sheet_names[sheet_num - 1]]
            
            # 데이터 분석
            text_content = ""
            non_empty_cells = 0
            total_cells = 0
            
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    total_cells += 1
                    if cell is not None:
                        non_empty_cells += 1
                        text_content += str(cell) + " "
            
            wb.close()
            
            # 메트릭 계산
            metrics = ContentMetrics()
            metrics.text_length = len(text_content.strip())
            metrics.word_count = len(text_content.split())
            metrics.sheet_count = 1
            metrics.non_empty_cell_ratio = non_empty_cells / total_cells if total_cells > 0 else 0
            
            # Excel은 기본적으로 text-based이지만, 빈 셀이 많으면 layout-based일 수 있음
            content_type = ContentType.TEXT_BASED
            confidence = 0.8
            reasoning = ["Excel 시트는 기본적으로 텍스트 기반"]
            
            if metrics.non_empty_cell_ratio < 0.1:
                content_type = ContentType.LAYOUT_BASED
                confidence = 0.6
                reasoning = ["빈 셀이 많아 레이아웃 중심으로 판단"]
            
            return ClassificationResult(
                content_type=content_type,
                confidence=confidence,
                reasoning=reasoning,
                metrics=metrics
            )
            
        except Exception as e:
            print(f"XLSX 시트 {sheet_num} 분석 오류: {e}")
            return ClassificationResult(
                content_type=ContentType.TEXT_BASED,
                confidence=0.5,
                reasoning=[f"분석 오류로 인한 기본값: {str(e)}"],
                metrics=ContentMetrics()
            )
    
    def _analyze_pdf_layout_complexity(self, text_blocks: Dict) -> float:
        """PDF 페이지의 레이아웃 복잡도 분석"""
        try:
            blocks = text_blocks.get("blocks", [])
            if not blocks:
                return 0.5  # 중간값
            
            # 텍스트 블록 수
            text_block_count = len([b for b in blocks if "lines" in b])
            
            # 폰트 다양성
            fonts = set()
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            fonts.add(span.get("font", ""))
            
            font_diversity = len(fonts)
            
            # 복잡도 점수 계산
            complexity = min((text_block_count / 10.0) + (font_diversity / 8.0), 1.0)
            
            return complexity
            
        except Exception:
            return 0.5
    
    def _classify_page_content(self, metrics: ContentMetrics, layout_complexity: float = 0.5) -> Tuple[ContentType, float, List[str]]:
        """페이지 콘텐츠 분류"""
        reasoning = []
        score = 0.5  # 기본값
        
        # 텍스트 길이 기준
        if metrics.text_length >= self.thresholds["min_text_length"]:
            score += 0.2
            reasoning.append(f"충분한 텍스트 ({metrics.text_length}자)")
        else:
            score -= 0.2
            reasoning.append(f"텍스트 부족 ({metrics.text_length}자)")
        
        # 이미지 비율 기준
        if metrics.image_count == 0:
            score += 0.15
            reasoning.append("이미지 없음")
        elif metrics.image_count <= 2:
            score += 0.05
            reasoning.append(f"적은 이미지 ({metrics.image_count}개)")
        else:
            score -= 0.15
            reasoning.append(f"많은 이미지 ({metrics.image_count}개)")
        
        # 텍스트 밀도 기준
        if hasattr(metrics, 'text_density') and metrics.text_density >= self.thresholds["min_text_density"]:
            score += 0.1
            reasoning.append(f"높은 텍스트 밀도 ({metrics.text_density:.2f})")
        
        # 레이아웃 복잡도 기준
        if layout_complexity <= self.thresholds["layout_complexity_threshold"]:
            score += 0.1
            reasoning.append(f"단순한 레이아웃 ({layout_complexity:.2f})")
        else:
            score -= 0.1
            reasoning.append(f"복잡한 레이아웃 ({layout_complexity:.2f})")
        
        # 최종 분류
        content_type = ContentType.TEXT_BASED if score >= 0.5 else ContentType.LAYOUT_BASED
        confidence = abs(score - 0.5) * 2  # 0.5에서 떨어진 정도를 신뢰도로 변환
        confidence = min(max(confidence, 0.1), 0.9)  # 0.1~0.9 범위로 제한
        
        return content_type, confidence, reasoning