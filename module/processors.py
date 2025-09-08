"""
파일 처리를 위한 프로세서 클래스들
"""

import os
import tempfile
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Dict, Any, Optional
from dotenv import load_dotenv
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from module.exceptions import ContentExtractionError, ProcessingError
from module.converters import ConverterFactory

# .env 파일 로드
load_dotenv()

logger = logging.getLogger(__name__)


class BaseProcessor(ABC):
    """파일 처리를 위한 추상 기본 클래스"""
    
    def __init__(self, azure_processor=None):
        self.azure_processor = azure_processor
    
    @abstractmethod
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """
        파일을 처리하여 JSON 형태의 청크 리스트를 반환합니다.
        
        Args:
            file_path: 처리할 파일 경로
            
        Returns:
            JSON 형태의 텍스트 청크 리스트
            
        Raises:
            ProcessingError: 처리 실패 시
        """
        pass
    
    def _create_text_chunk(self, text: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """텍스트 청크를 생성합니다."""
        return {
            "text_chunk_to_embed": text.strip(),
            "metadata": metadata
        }
    
    def _extract_section_title(self, text: str, previous_title: str = "") -> str:
        """텍스트에서 섹션 제목을 추출합니다."""
        # 간단한 제목 패턴 매칭
        lines = text.split('\n')
        for line in lines[:3]:  # 처음 3줄에서 제목 찾기
            line = line.strip()
            if line and len(line) < 100 and not line.startswith('-'):
                return line
        return previous_title


class TextBasedProcessor(BaseProcessor):
    """Text-based 파일 처리기"""
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """Text-based 파일을 처리하여 JSON 청크 리스트를 반환"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.docx':
                return self._process_docx(file_path)
            elif file_ext == '.pptx':
                return self._process_pptx(file_path)
            elif file_ext == '.xlsx':
                return self._process_xlsx(file_path)
            elif file_ext == '.pdf':
                return self._process_pdf_text(file_path)
            elif file_ext == '.txt':
                return self._process_txt(file_path)
            elif file_ext == '.md':
                return self._process_md(file_path)
            elif file_ext == '.csv':
                return self._process_csv(file_path)
            elif file_ext == '.scds':
                return self._process_scds(file_path)
            elif file_ext == '.xml':
                return self._process_xml(file_path)
            else:
                raise ContentExtractionError(
                    f"지원하지 않는 파일 형식: {file_ext}", 
                    file_path
                )
                
        except Exception as e:
            logger.error(f"Text-based 처리 오류: {e}")
            raise ProcessingError(f"Text-based 처리 실패: {str(e)}", file_path)
    
    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        DOCX 파일을 이중 경로 하이브리드 아키텍처로 처리
        
        1. 요소 단위 분석: 텍스트, 테이블, 이미지 추출
        2. 페이지 단위 분석: 전체 문서 Vision 분석 (PDF 변환 후)
        
        Args:
            file_path: DOCX 파일 경로
            
        Returns:
            이중 경로 분석 결과를 합성한 청크 리스트
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            logger.info(f"DOCX 이중 경로 하이브리드 처리 시작")
            
            # 경로 1: 요소 단위 분석
            element_analysis = self._analyze_docx_elements(doc, file_path)
            logger.info(f"DOCX 요소 단위 분석 완료: {len(element_analysis)}개 요소")
            
            # 경로 2: 페이지 단위 Vision 분석 (PDF 변환 후)
            page_analysis = self._analyze_docx_page_vision(file_path)
            logger.info(f"DOCX 페이지 단위 Vision 분석 완료")
            
            # 결과 합성: 페이지 단위 요약을 text_chunk_to_embed에, 요소 데이터를 metadata에
            hybrid_chunk = self._synthesize_docx_hybrid_result(file_path, element_analysis, page_analysis)
            
            return [hybrid_chunk]
            
        except Exception as e:
            logger.error(f"DOCX 파일 처리 오류: {e}")
            raise ContentExtractionError(f"DOCX 파일 처리 실패: {e}", file_path)
    
    def _process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        PPTX 파일을 이중 경로 하이브리드 아키텍처로 처리
        
        모든 슬라이드에 대해:
        1. 요소 단위 분석: 개별 도형, 텍스트, 표, 이미지 추출
        2. 페이지 단위 분석: 전체 슬라이드 Vision 분석
        
        Args:
            file_path: PPTX 파일 경로
            
        Returns:
            처리된 청크 리스트
        """
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            all_chunks = []
            
            logger.info(f"PPTX 이중 경로 하이브리드 처리 시작: {len(presentation.slides)}개 슬라이드")
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                logger.info(f"슬라이드 {slide_num} 이중 경로 분석 중...")
                
                # 이중 경로 하이브리드 처리
                slide_chunks = self._process_slide_hybrid(slide, slide_num, file_path)
                
                all_chunks.extend(slide_chunks)
                logger.info(f"슬라이드 {slide_num} 이중 경로 처리 완료: {len(slide_chunks)}개 청크")
            
            logger.info(f"PPTX 이중 경로 하이브리드 처리 완료: {len(all_chunks)}개 청크")
            return all_chunks
            
        except Exception as e:
            logger.error(f"PPTX 파일 처리 오류: {e}")
            raise ContentExtractionError(f"PPTX 파일 처리 실패: {e}", file_path)
    
    def _analyze_slide_content(self, slide) -> str:
        """
        개별 슬라이드의 콘텐츠를 분석하여 text_based 또는 layout_based 결정
        (기존 로직 - Vision-First로 대체됨)
        
        Args:
            slide: python-pptx의 slide 객체
            
        Returns:
            'text_based' 또는 'layout_based'
        """
        try:
            text_content = ""
            image_count = 0
            shape_count = len(slide.shapes)
            
            for shape in slide.shapes:
                # 텍스트 콘텐츠 수집
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text.strip() + "\n"
                
                # 이미지 개수 확인 (shape_type 13은 이미지)
                if shape.shape_type == 13:
                    image_count += 1
                
                # 차트, 다이어그램 등 복잡한 요소 확인
                if shape.shape_type in [17, 18, 19, 20]:  # 차트, SmartArt 등
                    image_count += 1
            
            # 텍스트 길이 계산
            text_length = len(text_content.strip())
            
            # 슬라이드 콘텐츠 타입 판별 로직
            if text_length > 100 and image_count < 3 and shape_count < 10:
                # 텍스트가 충분하고 이미지가 적고 복잡한 요소가 적음
                return 'text_based'
            elif text_length < 50 and image_count > 2:
                # 텍스트가 적고 이미지가 많음 (다이어그램, 차트 중심)
                return 'layout_based'
            elif shape_count > 15:
                # 복잡한 레이아웃과 많은 요소들
                return 'layout_based'
            else:
                # 기본적으로 텍스트 기반으로 처리
                return 'text_based'
                
        except Exception as e:
            logger.warning(f"슬라이드 콘텐츠 분석 오류: {e}")
            return 'text_based'  # 오류 시 안전하게 text_based로 분류
    
    # Vision-First 분류 전략은 이중 경로 하이브리드 아키텍처로 대체됨
    # 더 이상 text_based/layout_based 분류를 사용하지 않음
    
    def _process_slide_hybrid(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        슬라이드 이중 경로 하이브리드 처리
        
        경로 1: 요소 단위 분석 - 개별 도형, 텍스트, 표, 이미지 추출
        경로 2: 페이지 단위 분석 - 전체 슬라이드 Vision 분석
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            이중 경로 분석 결과를 합성한 청크 리스트
        """
        try:
            # Azure OpenAI 프로세서 초기화
            if not self.azure_processor:
                self.azure_processor = self._initialize_azure_processor()
            
            # 경로 1: 요소 단위 분석
            element_analysis = self._analyze_slide_elements(slide, slide_num, file_path)
            logger.info(f"슬라이드 {slide_num} 요소 단위 분석 완료: {len(element_analysis)}개 요소")
            
            # 경로 2: 페이지 단위 Vision 분석
            page_analysis = self._analyze_slide_page_vision(slide, slide_num, file_path)
            logger.info(f"슬라이드 {slide_num} 페이지 단위 Vision 분석 완료")
            
            # 결과 합성: 페이지 단위 요약을 text_chunk_to_embed에, 요소 데이터를 metadata에
            hybrid_chunk = self._synthesize_hybrid_result(
                slide_num, file_path, element_analysis, page_analysis
            )
            
            return [hybrid_chunk]
            
        except Exception as e:
            logger.error(f"슬라이드 {slide_num} 하이브리드 처리 오류: {e}")
            # 오류 시 기본 텍스트 기반 처리로 폴백
            return self._process_slide_text_based(slide, slide_num, file_path)
    
    def _analyze_slide_elements(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        슬라이드의 모든 요소를 개별적으로 분석
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            요소별 분석 결과 리스트
        """
        elements = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element_info = {
                    "shape_type": shape.shape_type,
                    "shape_name": shape.name,
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                }
                
                # 텍스트 상자 처리
                if hasattr(shape, "text") and shape.text.strip():
                    element_info.update({
                        "element_type": "text",
                        "content": shape.text.strip(),
                        "text_format": "plain"
                    })
                
                # 표 처리
                elif shape.shape_type == 19:  # 표
                    table_data = self._extract_table_data(shape)
                    element_info.update({
                        "element_type": "table",
                        "content": table_data,
                        "table_format": "markdown"
                    })
                
                # 이미지 처리
                elif shape.shape_type == 13:  # 이미지
                    image_description = self._analyze_image_with_vision(shape, slide_num, shape_idx)
                    element_info.update({
                        "element_type": "image",
                        "content": image_description,
                        "image_source": "pptx_slide"
                    })
                
                # 차트, SmartArt 등 복잡한 요소
                elif shape.shape_type in [17, 18, 20]:  # 차트, SmartArt 등
                    element_info.update({
                        "element_type": "complex_shape",
                        "content": f"복잡한 도형 요소 (타입: {shape.shape_type})",
                        "shape_category": "chart_or_smartart"
                    })
                
                # 기타 도형
                else:
                    element_info.update({
                        "element_type": "other_shape",
                        "content": f"기타 도형 요소 (타입: {shape.shape_type})"
                    })
                
                elements.append(element_info)
                
            except Exception as e:
                logger.warning(f"슬라이드 {slide_num} 요소 {shape_idx} 분석 오류: {e}")
                # 오류가 발생한 요소도 기본 정보로 저장
                elements.append({
                    "shape_type": getattr(shape, 'shape_type', 'unknown'),
                    "element_type": "error",
                    "content": f"요소 분석 실패: {str(e)}",
                    "error": str(e)
                })
        
        return elements
    
    def _extract_table_data(self, table_shape) -> str:
        """표 데이터를 마크다운 테이블로 변환"""
        try:
            table = table_shape.table
            markdown_table = []
            
            # 행과 열 수 확인
            if not table.rows or len(table.rows) == 0:
                return "[빈 표]"
            
            # 헤더 행 (첫 번째 행)
            header_row = []
            for cell in table.rows[0].cells:
                header_row.append(cell.text.strip() if cell.text else " ")
            markdown_table.append("| " + " | ".join(header_row) + " |")
            
            # 구분선
            markdown_table.append("| " + " | ".join(["---"] * len(header_row)) + " |")
            
            # 데이터 행 (두 번째 행부터)
            for i in range(1, len(table.rows)):
                row = table.rows[i]
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip() if cell.text else " ")
                markdown_table.append("| " + " | ".join(row_data) + " |")
            
            return "\n".join(markdown_table)
            
        except Exception as e:
            logger.warning(f"표 데이터 추출 오류: {e}")
            return f"[표 데이터 추출 실패: {str(e)}]"
    
    def _analyze_image_with_vision(self, image_shape, slide_num: int, shape_idx: int) -> str:
        """개별 이미지를 GPT-Vision으로 분석하여 짧은 설명 생성"""
        try:
            if not self.azure_processor:
                return f"[이미지: 추출 실패 - Azure OpenAI 프로세서 없음]"
            
            # 이미지 추출
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = self._extract_pptx_image(image_shape, slide_num, shape_idx)
                
                if not image_path or not os.path.exists(image_path):
                    return f"[이미지: 추출 실패 - 파일 없음]"
                
                # 이미지 설명용 프롬프트
                image_prompt = "이 이미지를 한 문장으로 간단하게 설명해주세요. 핵심 내용만 간결하게 표현해주세요."
                
                try:
                    description = self.azure_processor.image_to_text(image_path, image_prompt)
                    return description.strip()
                    
                except Exception as e:
                    logger.warning(f"이미지 Vision 분석 실패: {e}")
                    return f"[이미지: Vision 분석 실패 - {str(e)}]"
                    
        except Exception as e:
            logger.warning(f"이미지 분석 오류: {e}")
            return f"[이미지: 분석 오류 - {str(e)}]"
    
    def _analyze_slide_page_vision(self, slide, slide_num: int, file_path: str) -> str:
        """
        전체 슬라이드를 이미지로 변환하여 Vision으로 분석
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            전체 맥락에 대한 요약 텍스트
        """
        try:
            if not self.azure_processor:
                return f"슬라이드 {slide_num}: Azure OpenAI 프로세서 없음으로 인한 분석 제한"
            
            # 슬라이드를 이미지로 변환
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = self._convert_slide_to_image(slide, slide_num, temp_dir)
                
                if not image_path or not os.path.exists(image_path):
                    return f"슬라이드 {slide_num}: 이미지 변환 실패로 인한 분석 제한"
                
                # 전체 맥락 분석용 프롬프트
                context_prompt = """이 슬라이드의 핵심 주제와, 각 요소들 간의 관계를 종합적으로 설명해주세요.

다음 사항들을 포함하여 전체적인 맥락을 파악할 수 있도록 요약해주세요:
1. 슬라이드의 주요 주제나 목적
2. 제목과 내용의 관계
3. 텍스트, 이미지, 표 등 요소들의 배치와 연결성
4. 전체적인 메시지나 핵심 포인트

가독성 좋게 정리하여 슬라이드의 전체적인 의미를 이해할 수 있도록 해주세요."""
                
                try:
                    context_summary = self.azure_processor.image_to_text(image_path, context_prompt)
                    return context_summary.strip()
                    
                except Exception as e:
                    logger.warning(f"슬라이드 {slide_num} Vision 분석 실패: {e}")
                    return f"슬라이드 {slide_num}: Vision 분석 실패 - {str(e)}"
                    
        except Exception as e:
            logger.error(f"슬라이드 {slide_num} 페이지 Vision 분석 오류: {e}")
            return f"슬라이드 {slide_num}: 페이지 분석 오류 - {str(e)}"
    
    def _synthesize_hybrid_result(self, slide_num: int, file_path: str, 
                                 element_analysis: List[Dict], page_analysis: str) -> Dict[str, Any]:
        """
        이중 경로 분석 결과를 하나의 풍부한 JSON 객체로 합성
        
        Args:
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            element_analysis: 요소 단위 분석 결과
            page_analysis: 페이지 단위 Vision 분석 결과
            
        Returns:
            합성된 하이브리드 청크
        """
        # text_chunk_to_embed에는 전체 맥락 요약
        text_chunk = page_analysis
        
        # metadata에는 요소 단위 데이터와 추가 정보
        metadata = {
            "source_file": Path(file_path).name,
            "section_title": f"슬라이드 {slide_num}",
            "page_number": slide_num,
            "element_type": "hybrid_slide",
            "slide_number": slide_num,
            "file_type": "pptx",
            "processing_method": "dual_path_hybrid",
            "vision_analysis": True,
            "element_count": len(element_analysis),
            "elements": element_analysis,  # 모든 요소 분석 결과
            "page_context_summary": page_analysis,  # 페이지 단위 요약
            "architecture": "dual_path_hybrid"
        }
        
        return self._create_text_chunk(text_chunk, metadata)
    
    def _analyze_docx_elements(self, doc, file_path: str) -> List[Dict[str, Any]]:
        """DOCX 문서의 모든 요소를 개별적으로 분석"""
        elements = []
        current_section = "문서 시작"
        
        # 단락 처리
        for para_idx, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            
            # 스타일 정보 추출
            style_name = para.style.name if para.style else "Normal"
            
            # 제목인지 확인
            if style_name.startswith('Heading') or 'TOC' in style_name:
                current_section = text
            
            element_info = {
                "element_type": "text",
                "content": text,
                "style": style_name,
                "section": current_section,
                "paragraph_index": para_idx,
                "text_format": "plain"
            }
            elements.append(element_info)
        
        # 테이블 처리
        for table_idx, table in enumerate(doc.tables):
            try:
                table_text = self._table_to_markdown(table)
                if table_text.strip():
                    element_info = {
                        "element_type": "table",
                        "content": table_text,
                        "table_format": "markdown",
                        "table_index": table_idx,
                        "section": current_section
                    }
                    elements.append(element_info)
            except Exception as e:
                logger.warning(f"DOCX 테이블 {table_idx} 처리 오류: {e}")
                # 오류가 발생한 테이블도 기본 정보로 저장
                element_info = {
                    "element_type": "table_error",
                    "content": f"[테이블 처리 실패: {str(e)}]",
                    "table_index": table_idx,
                    "section": current_section,
                    "error": str(e)
                }
                elements.append(element_info)
        
        # 이미지 처리 (있는 경우)
        # DOCX의 이미지는 관계(relationship)를 통해 접근해야 함
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    element_info = {
                        "element_type": "image",
                        "content": f"[이미지: {rel.target_ref}]",
                        "image_source": "docx_relationship",
                        "target_ref": rel.target_ref
                    }
                    elements.append(element_info)
        except Exception as e:
            logger.warning(f"DOCX 이미지 추출 오류: {e}")
        
        return elements
    
    def _analyze_docx_page_vision(self, file_path: str) -> str:
        """DOCX 문서를 PDF로 변환하여 Vision으로 분석"""
        try:
            if not self.azure_processor:
                self.azure_processor = self._initialize_azure_processor()
            
            if not self.azure_processor:
                return f"DOCX: Azure OpenAI 프로세서 없음으로 인한 분석 제한"
            
            # DOCX를 PDF로 변환
            with tempfile.TemporaryDirectory() as temp_dir:
                pdf_path = self._convert_docx_to_pdf(file_path, temp_dir)
                
                if not pdf_path or not os.path.exists(pdf_path):
                    return f"DOCX: PDF 변환 실패로 인한 분석 제한"
                
                # 전체 문서 맥락 분석용 프롬프트
                context_prompt = """이 문서의 핵심 주제와, 각 요소들 간의 관계를 종합적으로 설명해주세요.

다음 사항들을 포함하여 전체적인 맥락을 파악할 수 있도록 요약해주세요:
1. 문서의 주요 주제나 목적
2. 제목과 내용의 구조적 관계
3. 텍스트, 테이블, 이미지 등 요소들의 배치와 연결성
4. 전체적인 메시지나 핵심 포인트

가독성 좋게 정리하여 문서의 전체적인 의미를 이해할 수 있도록 해주세요."""
                
                try:
                    context_summary = self.azure_processor.image_to_text(pdf_path, context_prompt)
                    return context_summary.strip()
                    
                except Exception as e:
                    logger.warning(f"DOCX Vision 분석 실패: {e}")
                    return f"DOCX: Vision 분석 실패 - {str(e)}"
                    
        except Exception as e:
            logger.error(f"DOCX 페이지 Vision 분석 오류: {e}")
            return f"DOCX: 페이지 분석 오류 - {str(e)}"
    
    def _convert_docx_to_pdf(self, docx_path: str, temp_dir: str) -> str:
        """DOCX를 PDF로 변환"""
        try:
            from module.converters import ConverterFactory
            
            converter = ConverterFactory.get_best_converter("docx", "pdf")
            if not converter.is_available():
                logger.warning("사용 가능한 DOCX-PDF 변환기가 없음")
                return None
            
            pdf_path = os.path.join(temp_dir, "document.pdf")
            if converter.convert(docx_path, pdf_path):
                return pdf_path
            else:
                return None
                
        except Exception as e:
            logger.warning(f"DOCX-PDF 변환 오류: {e}")
            return None
    
    def _synthesize_docx_hybrid_result(self, file_path: str, element_analysis: List[Dict], page_analysis: str) -> Dict[str, Any]:
        """DOCX 이중 경로 분석 결과를 하나의 풍부한 JSON 객체로 합성"""
        # text_chunk_to_embed에는 전체 맥락 요약
        text_chunk = page_analysis
        
        # metadata에는 요소 단위 데이터와 추가 정보
        metadata = {
            "source_file": Path(file_path).name,
            "section_title": "DOCX 문서 전체",
            "page_number": 1,
            "element_type": "hybrid_document",
            "file_type": "docX",
            "processing_method": "dual_path_hybrid",
            "vision_analysis": True,
            "element_count": len(element_analysis),
            "elements": element_analysis,  # 모든 요소 분석 결과
            "page_context_summary": page_analysis,  # 페이지 단위 요약
            "architecture": "dual_path_hybrid"
        }
        
        return self._create_text_chunk(text_chunk, metadata)
    
    def _process_pdf_page_hybrid(self, page, page_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        PDF 페이지 이중 경로 하이브리드 처리
        
        경로 1: 요소 단위 분석 - 텍스트, 테이블, 이미지 추출
        경로 2: 페이지 단위 분석 - 전체 페이지 Vision 분석
        
        Args:
            page: PDF 페이지 객체
            page_num: 페이지 번호
            file_path: 원본 파일 경로
            
        Returns:
            이중 경로 분석 결과를 합성한 청크 리스트
        """
        try:
            # Azure OpenAI 프로세서 초기화
            if not self.azure_processor:
                self.azure_processor = self._initialize_azure_processor()
            
            # 경로 1: 요소 단위 분석
            element_analysis = self._analyze_pdf_page_elements(page, page_num, file_path)
            logger.info(f"PDF 페이지 {page_num} 요소 단위 분석 완료: {len(element_analysis)}개 요소")
            
            # 경로 2: 페이지 단위 Vision 분석
            page_analysis = self._analyze_pdf_page_vision(page, page_num, file_path)
            logger.info(f"PDF 페이지 {page_num} Vision 분석 완료")
            
            # 결과 합성: 페이지 단위 요약을 text_chunk_to_embed에, 요소 데이터를 metadata에
            hybrid_chunk = self._synthesize_pdf_page_hybrid_result(
                page_num, file_path, element_analysis, page_analysis
            )
            
            return [hybrid_chunk]
            
        except Exception as e:
            logger.error(f"PDF 페이지 {page_num} 하이브리드 처리 오류: {e}")
            # 오류 시 기본 텍스트 기반 처리로 폴백
            return self._process_pdf_page_text_based(page, page_num, file_path)
    
    def _analyze_pdf_page_elements(self, page, page_num: int, file_path: str) -> List[Dict[str, Any]]:
        """PDF 페이지의 모든 요소를 개별적으로 분석"""
        elements = []
        
        # 텍스트 추출
        text = page.get_text()
        if text.strip():
            element_info = {
                "element_type": "text",
                "content": text.strip(),
                "text_format": "plain",
                "page_number": page_num
            }
            elements.append(element_info)
        
        # 테이블 추출 (PyMuPDF의 테이블 감지 기능 사용)
        try:
            # PyMuPDF 버전에 따라 다른 메서드 사용
            if hasattr(page, 'get_tables'):
                tables = page.get_tables()
                for table_idx, table in enumerate(tables):
                    if table:
                        table_text = self._table_to_markdown(table)
                        if table_text.strip():
                            element_info = {
                                "element_type": "table",
                                "content": table_text,
                                "table_format": "markdown",
                                "table_index": table_idx,
                                "page_number": page_num
                            }
                            elements.append(element_info)
            else:
                # 테이블 감지 기능이 없는 경우 텍스트에서 테이블 패턴 찾기
                text = page.get_text()
                lines = text.split('\n')
                table_lines = []
                in_table = False
                
                for line in lines:
                    if '\t' in line or '  ' in line:  # 탭이나 여러 공백으로 구분된 경우
                        if not in_table:
                            in_table = True
                        table_lines.append(line)
                    elif in_table:
                        in_table = False
                        if table_lines:
                            # 간단한 테이블로 처리
                            element_info = {
                                "element_type": "table_detected",
                                "content": "\n".join(table_lines),
                                "table_format": "text_pattern",
                                "page_number": page_num
                            }
                            elements.append(element_info)
                            table_lines = []
                
                # 마지막 테이블 처리
                if table_lines:
                    element_info = {
                        "element_type": "table_detected",
                        "content": "\n".join(table_lines),
                        "table_format": "text_pattern",
                        "page_number": page_num
                    }
                    elements.append(element_info)
                    
        except Exception as e:
            logger.warning(f"PDF 페이지 {page_num} 테이블 추출 오류: {e}")
        
        # 이미지 추출
        try:
            image_list = page.get_images()
            for img_idx, img in enumerate(image_list):
                element_info = {
                    "element_type": "image",
                    "content": f"[이미지: {img[0]}]",
                    "image_source": "pdf_page",
                    "image_index": img_idx,
                    "page_number": page_num
                }
                elements.append(element_info)
        except Exception as e:
            logger.warning(f"PDF 페이지 {page_num} 이미지 추출 오류: {e}")
        
        return elements
    
    def _analyze_pdf_page_vision(self, page, page_num: int, file_path: str) -> str:
        """PDF 페이지를 이미지로 변환하여 Vision으로 분석"""
        try:
            if not self.azure_processor:
                return f"PDF 페이지 {page_num}: Azure OpenAI 프로세서 없음으로 인한 분석 제한"
            
            # PDF 페이지를 이미지로 변환
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = self._convert_pdf_page_to_image(page, page_num, temp_dir)
                
                if not image_path or not os.path.exists(image_path):
                    return f"PDF 페이지 {page_num}: 이미지 변환 실패로 인한 분석 제한"
                
                # 전체 페이지 맥락 분석용 프롬프트
                context_prompt = """이 PDF 페이지의 핵심 주제와, 각 요소들 간의 관계를 종합적으로 설명해주세요.

다음 사항들을 포함하여 전체적인 맥락을 파악할 수 있도록 요약해주세요:
1. 페이지의 주요 주제나 목적
2. 제목과 내용의 구조적 관계
3. 텍스트, 테이블, 이미지 등 요소들의 배치와 연결성
4. 전체적인 메시지나 핵심 포인트

가독성 좋게 정리하여 페이지의 전체적인 의미를 이해할 수 있도록 해주세요."""
                
                try:
                    context_summary = self.azure_processor.image_to_text(image_path, context_prompt)
                    return context_summary.strip()
                    
                except Exception as e:
                    logger.warning(f"PDF 페이지 {page_num} Vision 분석 실패: {e}")
                    return f"PDF 페이지 {page_num}: Vision 분석 실패 - {str(e)}"
                    
        except Exception as e:
            logger.error(f"PDF 페이지 {page_num} Vision 분석 오류: {e}")
            return f"PDF 페이지 {page_num}: Vision 분석 오류 - {str(e)}"
    
    def _convert_pdf_page_to_image(self, page, page_num: int, temp_dir: str) -> str:
        """PDF 페이지를 이미지로 변환"""
        try:
            import fitz  # PyMuPDF
            
            # 페이지를 이미지로 렌더링
            mat = fitz.Matrix(2.0, 2.0)  # 2배 확대
            pix = page.get_pixmap(matrix=mat)
            
            image_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(image_path)
            
            return image_path
            
        except Exception as e:
            logger.warning(f"PDF 페이지 {page_num} 이미지 변환 오류: {e}")
            return None
    
    def _synthesize_pdf_page_hybrid_result(self, page_num: int, file_path: str, 
                                          element_analysis: List[Dict], page_analysis: str) -> Dict[str, Any]:
        """PDF 페이지 이중 경로 분석 결과를 하나의 풍부한 JSON 객체로 합성"""
        # text_chunk_to_embed에는 전체 맥락 요약
        text_chunk = page_analysis
        
        # metadata에는 요소 단위 데이터와 추가 정보
        metadata = {
            "source_file": Path(file_path).name,
            "section_title": f"PDF 페이지 {page_num}",
            "page_number": page_num,
            "element_type": "hybrid_pdf_page",
            "file_type": "pdf",
            "processing_method": "dual_path_hybrid",
            "vision_analysis": True,
            "element_count": len(element_analysis),
            "elements": element_analysis,  # 모든 요소 분석 결과
            "page_context_summary": page_analysis,  # 페이지 단위 요약
            "architecture": "dual_path_hybrid"
        }
        
        return self._create_text_chunk(text_chunk, metadata)
    
    def _synthesize_simple_result(self, file_path: str, element_analysis: List[Dict], file_type: str) -> Dict[str, Any]:
        """
        단순 변환 방식 결과를 하나의 JSON 객체로 합성
        
        text_chunk_to_embed에는 빈 문자열, 요소 데이터는 metadata에
        
        Args:
            file_path: 원본 파일 경로
            element_analysis: 요소 단위 분석 결과
            file_type: 파일 타입
            
        Returns:
            단순 변환 결과 청크
        """
        # text_chunk_to_embed에는 빈 문자열 (Vision 분석 없음)
        text_chunk = ""
        
        # metadata에는 요소 단위 데이터와 추가 정보
        metadata = {
            "source_file": Path(file_path).name,
            "section_title": f"{file_type.upper()} 파일 전체",
            "page_number": 1,
            "element_type": "simple_conversion",
            "file_type": file_type,
            "processing_method": "simple_conversion",
            "vision_analysis": False,
            "element_count": len(element_analysis),
            "elements": element_analysis,  # 모든 요소 분석 결과
            "architecture": "simple_conversion"
        }
        
        return self._create_text_chunk(text_chunk, metadata)
    
    def _process_txt(self, file_path: str) -> List[Dict[str, Any]]:
        """
        TXT 파일을 단순 변환 방식으로 처리
        
        요소 단위 분석만 수행 (Vision 분석 없음)
        text_chunk_to_embed에는 빈 문자열
        
        Args:
            file_path: TXT 파일 경로
            
        Returns:
            요소 단위 분석 결과를 metadata에 포함한 청크 리스트
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            logger.info(f"TXT 단순 변환 처리 시작")
            
            # 텍스트를 줄 단위로 분석
            lines = content.split('\n')
            all_elements = []
            
            for line_idx, line in enumerate(lines):
                line = line.strip()
                if line:  # 빈 줄 제외
                    element_info = {
                        "element_type": "text_line",
                        "content": line,
                        "line_number": line_idx + 1,
                        "text_format": "plain"
                    }
                    all_elements.append(element_info)
            
            # 결과 합성: text_chunk_to_embed는 빈 문자열, 요소 데이터는 metadata에
            simple_chunk = self._synthesize_simple_result(file_path, all_elements, "txt")
            
            logger.info(f"TXT 단순 변환 처리 완료: {len(all_elements)}개 요소")
            return [simple_chunk]
            
        except Exception as e:
            logger.error(f"TXT 파일 처리 오류: {e}")
            raise ContentExtractionError(f"TXT 파일 처리 실패: {e}", file_path)
    
    def _process_md(self, file_path: str) -> List[Dict[str, Any]]:
        """
        MD 파일을 단순 변환 방식으로 처리
        
        요소 단위 분석만 수행 (Vision 분석 없음)
        text_chunk_to_embed에는 빈 문자열
        
        Args:
            file_path: MD 파일 경로
            
        Returns:
            요소 단위 분석 결과를 metadata에 포함한 청크 리스트
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            logger.info(f"MD 단순 변환 처리 시작")
            
            # 마크다운을 구조별로 분석
            lines = content.split('\n')
            all_elements = []
            current_section = "문서 시작"
            
            for line_idx, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                # 제목 감지
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    current_section = line.lstrip('#').strip()
                    
                    element_info = {
                        "element_type": "heading",
                        "content": line,
                        "heading_level": level,
                        "section": current_section,
                        "line_number": line_idx + 1,
                        "text_format": "markdown"
                    }
                    all_elements.append(element_info)
                
                # 리스트 감지
                elif line.startswith('- ') or line.startswith('* '):
                    element_info = {
                        "element_type": "list_item",
                        "content": line,
                        "section": current_section,
                        "line_number": line_idx + 1,
                        "text_format": "markdown"
                    }
                    all_elements.append(element_info)
                
                # 코드 블록 감지
                elif line.startswith('```'):
                    element_info = {
                        "element_type": "code_block",
                        "content": line,
                        "section": current_section,
                        "line_number": line_idx + 1,
                        "text_format": "markdown"
                    }
                    all_elements.append(element_info)
                
                # 일반 텍스트
                else:
                    element_info = {
                        "element_type": "text",
                        "content": line,
                        "section": current_section,
                        "line_number": line_idx + 1,
                        "text_format": "markdown"
                    }
                    all_elements.append(element_info)
            
            # 결과 합성: text_chunk_to_embed는 빈 문자열, 요소 데이터는 metadata에
            simple_chunk = self._synthesize_simple_result(file_path, all_elements, "md")
            
            logger.info(f"MD 단순 변환 처리 완료: {len(all_elements)}개 요소")
            return [simple_chunk]
            
        except Exception as e:
            logger.error(f"MD 파일 처리 오류: {e}")
            raise ContentExtractionError(f"MD 파일 처리 실패: {e}", file_path)
    
    def _process_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """
        CSV 파일을 단순 변환 방식으로 처리
        
        요소 단위 분석만 수행 (Vision 분석 없음)
        text_chunk_to_embed에는 빈 문자열
        
        Args:
            file_path: CSV 파일 경로
            
        Returns:
            요소 단위 분석 결과를 metadata에 포함한 청크 리스트
        """
        try:
            import csv
            
            logger.info(f"CSV 단순 변환 처리 시작")
            
            all_elements = []
            
            with open(file_path, 'r', encoding='utf-8') as f:
                csv_reader = csv.reader(f)
                rows = list(csv_reader)
                
                if rows:
                    # 헤더 행
                    header_row = rows[0]
                    header_element = {
                        "element_type": "header",
                        "content": header_row,
                        "row_type": "header",
                        "column_count": len(header_row),
                        "text_format": "csv"
                    }
                    all_elements.append(header_element)
                    
                    # 데이터 행들
                    for row_idx, row in enumerate(rows[1:], 1):
                        if row:  # 빈 행 제외
                            element_info = {
                                "element_type": "data_row",
                                "content": row,
                                "row_type": "data",
                                "row_number": row_idx,
                                "column_count": len(row),
                                "text_format": "csv"
                            }
                            all_elements.append(element_info)
            
            # 결과 합성: text_chunk_to_embed는 빈 문자열, 요소 데이터는 metadata에
            simple_chunk = self._synthesize_simple_result(file_path, all_elements, "csv")
            
            logger.info(f"CSV 단순 변환 처리 완료: {len(all_elements)}개 요소")
            return [simple_chunk]
            
        except Exception as e:
            logger.error(f"CSV 파일 처리 오류: {e}")
            raise ContentExtractionError(f"CSV 파일 처리 실패: {e}", file_path)
    
    def _process_scds(self, file_path: str) -> List[Dict[str, Any]]:
        """
        SCDS 바이너리 파일을 단순 변환 방식으로 처리
        
        바이너리 파일의 기본 정보와 구조를 분석하여 메타데이터로 저장
        
        Args:
            file_path: SCDS 파일 경로
            
        Returns:
            바이너리 분석 결과를 metadata에 포함한 청크 리스트
        """
        try:
            logger.info(f"SCDS 바이너리 파일 처리 시작")
            
            all_elements = []
            
            with open(file_path, 'rb') as f:
                # 파일 헤더 분석
                header = f.read(16)
                header_hex = ' '.join(f'{b:02x}' for b in header)
                header_ascii = ''.join(chr(b) if 32 <= b <= 126 else '.' for b in header)
                
                # 헤더 정보
                header_element = {
                    "element_type": "binary_header",
                    "content": f"SCDS Header: {header_ascii}",
                    "hex_data": header_hex,
                    "ascii_data": header_ascii,
                    "data_format": "binary"
                }
                all_elements.append(header_element)
                
                # 파일 크기
                file_size = os.path.getsize(file_path)
                size_element = {
                    "element_type": "file_info",
                    "content": f"File Size: {file_size:,} bytes ({file_size/1024:.1f} KB)",
                    "file_size_bytes": file_size,
                    "file_size_kb": file_size / 1024,
                    "data_format": "metadata"
                }
                all_elements.append(size_element)
                
                # 바이너리 패턴 분석
                f.seek(0)
                content = f.read()
                
                # 텍스트 문자열 추출 시도
                text_strings = []
                current_string = ""
                
                for byte in content:
                    if 32 <= byte <= 126:  # 출력 가능한 ASCII 문자
                        current_string += chr(byte)
                    else:
                        if len(current_string) >= 3:  # 3자 이상의 문자열만 저장
                            text_strings.append(current_string)
                        current_string = ""
                
                # 마지막 문자열 처리
                if len(current_string) >= 3:
                    text_strings.append(current_string)
                
                # 의미 있는 문자열들 저장
                for i, text in enumerate(text_strings[:20]):  # 처음 20개만
                    if text.strip():
                        text_element = {
                            "element_type": "extracted_text",
                            "content": text.strip(),
                            "text_index": i,
                            "text_length": len(text),
                            "data_format": "text"
                        }
                        all_elements.append(text_element)
                
                # 바이너리 통계
                binary_stats = {
                    "element_type": "binary_statistics",
                    "content": f"Binary Analysis: {len(text_strings)} text strings found",
                    "total_strings": len(text_strings),
                    "avg_string_length": sum(len(s) for s in text_strings) / max(len(text_strings), 1),
                    "data_format": "statistics"
                }
                all_elements.append(binary_stats)
            
            # 결과 합성: text_chunk_to_embed는 빈 문자열, 요소 데이터는 metadata에
            simple_chunk = self._synthesize_simple_result(file_path, all_elements, "scds")
            
            logger.info(f"SCDS 바이너리 파일 처리 완료: {len(all_elements)}개 요소")
            return [simple_chunk]
            
        except Exception as e:
            logger.error(f"SCDS 파일 처리 오류: {e}")
            raise ContentExtractionError(f"SCDS 파일 처리 실패: {e}", file_path)
    
    def _process_pdf_page_text_based(self, page, page_num: int, file_path: str) -> List[Dict[str, Any]]:
        """PDF 페이지 기본 텍스트 기반 처리 (폴백용)"""
        try:
            text = page.get_text()
            if text.strip():
                metadata = {
                    "source_file": Path(file_path).name,
                    "section_title": f"PDF 페이지 {page_num}",
                    "page_number": page_num,
                    "element_type": "text",
                    "file_type": "pdf",
                    "processing_method": "text_based_fallback"
                }
                
                chunk = self._create_text_chunk(text.strip(), metadata)
                return [chunk]
            else:
                return []
                
        except Exception as e:
            logger.error(f"PDF 페이지 {page_num} 텍스트 기반 처리 오류: {e}")
            return []
    
    def _initialize_azure_processor(self):
        """환경 변수에서 Azure OpenAI 설정을 읽어 프로세서를 초기화"""
        try:
            from module.image_to_text import AzureOpenAIImageProcessor
            
            # 환경 변수에서 설정 읽기
            api_key = os.getenv("AZURE_OPENAI_API_KEY")
            endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
            deployment_name = os.getenv("AZURE_OPENAI_DEPLOYMENT_NAME", "gpt-4-vision")
            api_version = os.getenv("AZURE_OPENAI_API_VERSION", "2024-10-21")
            
            if not api_key or not endpoint:
                logger.warning("Azure OpenAI 환경 변수가 설정되지 않음")
                return None
            
            # Azure OpenAI 프로세서 초기화
            azure_processor = AzureOpenAIImageProcessor(api_key, endpoint, deployment_name)
            logger.info("Azure OpenAI 프로세서 자동 초기화 성공")
            return azure_processor
            
        except Exception as e:
            logger.error(f"Azure OpenAI 프로세서 초기화 실패: {e}")
            return None
    
    def _process_slide_text_based(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        텍스트 기반 슬라이드 처리: 개별 요소들을 추출하여 청크 생성
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            텍스트 기반 처리로 생성된 청크 리스트
        """
        chunks = []
        slide_texts = []
        
        # 슬라이드의 모든 텍스트 수집
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        
        if slide_texts:
            # 슬라이드 제목과 내용을 구분
            slide_title = slide_texts[0] if slide_texts else f"슬라이드 {slide_num}"
            slide_content = "\n".join(slide_texts[1:]) if len(slide_texts) > 1 else ""
            
            # 슬라이드 제목 청크
            title_metadata = {
                "source_file": Path(file_path).name,
                "section_title": f"슬라이드 {slide_num}",
                "page_number": slide_num,
                "element_type": "slide_title",
                "slide_number": slide_num,
                "file_type": "pptx",
                "processing_method": "text_based"
            }
            
            title_chunk = self._create_text_chunk(slide_title, title_metadata)
            chunks.append(title_chunk)
            
            # 슬라이드 내용 청크
            if slide_content:
                content_metadata = {
                    "source_file": Path(file_path).name,
                    "section_title": slide_title,
                    "page_number": slide_num,
                    "element_type": "slide_content",
                    "slide_number": slide_num,
                    "file_type": "pptx",
                    "processing_method": "text_based"
                }
                
                content_chunk = self._create_text_chunk(slide_content, content_metadata)
                chunks.append(content_chunk)
        
        # 이미지가 있는 경우 Vision 기반 처리로 전환
        image_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
        if image_shapes and self.azure_processor:
            logger.info(f"슬라이드 {slide_num}: 이미지가 감지되어 Vision 기반 처리로 전환")
            
            # 슬라이드를 이미지로 변환하여 Vision 분석
            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    image_path = self._convert_slide_to_image(slide, slide_num, temp_dir)
                    
                    if image_path and os.path.exists(image_path):
                        # 이미지 내용 분석용 프롬프트
                        image_analysis_prompt = """이 슬라이드에 포함된 이미지들을 분석하고 설명해주세요.

다음 사항들을 포함하여 상세하게 설명해주세요:
1. 이미지의 종류와 내용 (아이콘, 사진, 다이어그램 등)
2. 이미지가 전달하는 의미나 메시지
3. 텍스트와의 관계 및 전체적인 맥락
4. 중요한 시각적 요소나 특징

가독성 좋게 정리하여 이미지의 내용을 명확하게 이해할 수 있도록 해주세요."""
                        
                        try:
                            image_description = self.azure_processor.image_to_text(image_path, image_analysis_prompt)
                            
                            # Vision 기반 이미지 분석 결과를 청크로 저장
                            image_metadata = {
                                "source_file": Path(file_path).name,
                                "section_title": f"슬라이드 {slide_num}",
                                "page_number": slide_num,
                                "element_type": "vision_analyzed_image",
                                "slide_number": slide_num,
                                "file_type": "pptx",
                                "processing_method": "vision_enhanced_text",
                                "image_source": "pptx_slide",
                                "vision_analysis": True,
                                "gpt_vision_prompt": image_analysis_prompt
                            }
                            
                            image_chunk = self._create_text_chunk(image_description, image_metadata)
                            chunks.append(image_chunk)
                            
                            logger.info(f"슬라이드 {slide_num} 이미지 Vision 분석 완료")
                            
                        except Exception as e:
                            logger.warning(f"슬라이드 {slide_num} 이미지 Vision 분석 실패: {e}")
                            # Vision 분석 실패 시 기본 이미지 정보만 저장
                            self._add_basic_image_info(chunks, slide_num, file_path, image_shapes)
                            
                    else:
                        logger.warning(f"슬라이드 {slide_num}: 이미지 변환 실패, 기본 정보만 저장")
                        self._add_basic_image_info(chunks, slide_num, file_path, image_shapes)
                        
                except Exception as e:
                    logger.warning(f"슬라이드 {slide_num} Vision 처리 실패: {e}, 기본 정보만 저장")
                    self._add_basic_image_info(chunks, slide_num, file_path, image_shapes)
        elif image_shapes:
            # Azure OpenAI 프로세서가 없는 경우 기본 이미지 정보만 저장
            self._add_basic_image_info(chunks, slide_num, file_path, image_shapes)
        
        logger.info(f"슬라이드 {slide_num} 텍스트 기반 처리 완료: {len(chunks)}개 청크")
        return chunks
    
    def _add_basic_image_info(self, chunks: List[Dict[str, Any]], slide_num: int, file_path: str, image_shapes: List):
        """기본 이미지 정보를 청크에 추가"""
        for i, shape in enumerate(image_shapes):
            try:
                image_path = self._extract_pptx_image(shape, slide_num, i)
                image_metadata = {
                    "source_file": Path(file_path).name,
                    "section_title": f"슬라이드 {slide_num}",
                    "page_number": slide_num,
                    "element_type": "image",
                    "slide_number": slide_num,
                    "file_type": "pptx",
                    "processing_method": "text_based",
                    "image_source": "pptx_slide",
                    "vision_analysis": False
                }
                
                image_chunk = self._create_text_chunk(f"[이미지: {image_path}]", image_metadata)
                chunks.append(image_chunk)
                
            except Exception as e:
                logger.warning(f"이미지 {i} 추출 오류: {e}")
                # 이미지 추출 실패 시에도 기본 정보는 저장
                image_metadata = {
                    "source_file": Path(file_path).name,
                    "section_title": f"슬라이드 {slide_num}",
                    "page_number": slide_num,
                    "element_type": "image",
                    "slide_number": slide_num,
                    "file_type": "pptx",
                    "processing_method": "text_based",
                    "image_source": "pptx_slide",
                    "vision_analysis": False,
                    "extraction_error": str(e)
                }
                
                image_chunk = self._create_text_chunk(f"[이미지: 추출 실패 - {str(e)}]", image_metadata)
                chunks.append(image_chunk)
    
    def _process_slide_layout_based(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        레이아웃 기반 슬라이드 처리: 이미지로 변환 후 GPT-Vision 처리
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            레이아웃 기반 처리로 생성된 청크 리스트
        """
        chunks = []
        
        try:
            if not self.azure_processor:
                logger.warning(f"슬라이드 {slide_num}: Azure OpenAI 프로세서가 없어 텍스트 기반 처리로 대체")
                return self._process_slide_text_based(slide, slide_num, file_path)
            
            # 슬라이드를 이미지로 변환
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = self._convert_slide_to_image(slide, slide_num, temp_dir)
                
                # GPT Vision으로 이미지 처리
                prompt = "이 슬라이드의 모든 텍스트 내용을 추출하고, 표나 이미지가 있다면 설명해주세요. 구조화된 형태로 정리해주세요."
                text_content = self.azure_processor.image_to_text(image_path, prompt)
                
                # 결과를 청크로 저장
                metadata = {
                    "source_file": Path(file_path).name,
                    "section_title": f"슬라이드 {slide_num}",
                    "page_number": slide_num,
                    "element_type": "image_processed",
                    "slide_number": slide_num,
                    "file_type": "pptx",
                    "processing_method": "layout_based",
                    "gpt_vision_prompt": prompt
                }
                
                chunk = self._create_text_chunk(text_content, metadata)
                chunks.append(chunk)
                
                logger.info(f"슬라이드 {slide_num} 레이아웃 기반 처리 완료: 1개 청크")
                return chunks
                
        except Exception as e:
            logger.error(f"슬라이드 {slide_num} 레이아웃 기반 처리 오류: {e}")
            logger.info(f"슬라이드 {slide_num}: 레이아웃 기반 처리 실패, 텍스트 기반 처리로 대체")
            return self._process_slide_text_based(slide, slide_num, file_path)
    
    def _process_slide_layout_based_vision(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """
        Vision-First 레이아웃 기반 슬라이드 처리: GPT-Vision으로 상세 내용 추출
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            file_path: 원본 파일 경로
            
        Returns:
            Vision 기반 처리로 생성된 청크 리스트
        """
        chunks = []
        
        try:
            # Azure OpenAI 프로세서가 없으면 자동으로 초기화 시도
            if not self.azure_processor:
                self.azure_processor = self._initialize_azure_processor()
                
            if not self.azure_processor:
                logger.warning(f"슬라이드 {slide_num}: Azure OpenAI 프로세서 초기화 실패, 텍스트 기반 처리로 대체")
                return self._process_slide_text_based(slide, slide_num, file_path)
            
            # 슬라이드를 이미지로 변환
            with tempfile.TemporaryDirectory() as temp_dir:
                image_path = self._convert_slide_to_image(slide, slide_num, temp_dir)
                
                if not image_path or not os.path.exists(image_path):
                    logger.warning(f"슬라이드 {slide_num}: 이미지 변환 실패, 텍스트 기반 처리로 대체")
                    return self._process_slide_text_based(slide, slide_num, file_path)
                
                # 내용 추출용 프롬프트로 GPT-Vision 호출
                content_extraction_prompt = """이 슬라이드의 모든 내용을 상세하게 분석하고 추출해주세요.

다음 사항들을 포함하여 구조화된 형태로 정리해주세요:
1. 모든 텍스트 내용 (제목, 본문, 라벨 등)
2. 표나 리스트의 구조와 내용
3. 이미지나 아이콘의 의미와 설명
4. 요소들 간의 관계와 레이아웃 정보
5. 중요한 정보나 핵심 포인트

가독성 좋게 정리하여 사용자가 슬라이드의 내용을 완전히 이해할 수 있도록 해주세요."""
                
                try:
                    text_content = self.azure_processor.image_to_text(image_path, content_extraction_prompt)
                    
                    # 결과를 청크로 저장
                    metadata = {
                        "source_file": Path(file_path).name,
                        "section_title": f"슬라이드 {slide_num}",
                        "slide_number": slide_num,
                        "element_type": "vision_processed",
                        "page_number": slide_num,
                        "file_type": "pptx",
                        "processing_method": "vision_first_layout",
                        "gpt_vision_prompt": content_extraction_prompt,
                        "vision_classification": "layout_based"
                    }
                    
                    chunk = self._create_text_chunk(text_content, metadata)
                    chunks.append(chunk)
                    
                    logger.info(f"슬라이드 {slide_num} Vision-First 레이아웃 처리 완료: 1개 청크")
                    return chunks
                    
                except Exception as e:
                    logger.error(f"슬라이드 {slide_num} Vision 내용 추출 실패: {e}")
                    logger.info(f"슬라이드 {slide_num}: Vision 내용 추출 실패, 텍스트 기반 처리로 대체")
                    return self._process_slide_text_based(slide, slide_num, file_path)
                    
        except Exception as e:
            logger.error(f"슬라이드 {slide_num} Vision-First 레이아웃 처리 오류: {e}")
            logger.info(f"슬라이드 {slide_num}: Vision-First 레이아웃 처리 실패, 텍스트 기반 처리로 대체")
            return self._process_slide_text_based(slide, slide_num, file_path)
    
    def _convert_slide_to_image(self, slide, slide_num: int, temp_dir: str) -> str:
        """
        슬라이드를 이미지로 변환 (LibreOffice를 통한 간접 변환)
        
        Args:
            slide: 슬라이드 객체
            slide_num: 슬라이드 번호
            temp_dir: 임시 디렉토리 경로
            
        Returns:
            생성된 이미지 파일 경로
        """
        try:
            # python-pptx는 직접 이미지 변환을 지원하지 않으므로
            # 대안 방법들을 시도
            
            # 방법 1: 슬라이드 텍스트를 임시 PPTX로 저장 후 LibreOffice로 변환
            temp_pptx_path = self._create_temp_pptx_with_slide(slide, slide_num, temp_dir)
            
            if temp_pptx_path and os.path.exists(temp_pptx_path):
                # LibreOffice를 사용하여 PPTX -> PDF -> 이미지 변환
                pdf_path = self._convert_pptx_to_pdf(temp_pptx_path, temp_dir)
                if pdf_path and os.path.exists(pdf_path):
                    image_path = self._convert_pdf_page_to_image(pdf_path, 1, slide_num, temp_dir)
                    if image_path and os.path.exists(image_path):
                        logger.info(f"슬라이드 {slide_num} 이미지 변환 성공: {image_path}")
                        return image_path
            
            # 방법 2: 텍스트 기반 처리로 대체
            logger.warning(f"슬라이드 {slide_num}: 이미지 변환 실패, 텍스트 기반 처리로 대체")
            raise NotImplementedError("슬라이드 이미지 변환에 필요한 의존성이 없음")
            
        except Exception as e:
            logger.error(f"슬라이드 {slide_num} 이미지 변환 오류: {e}")
            raise e
    
    def _create_temp_pptx_with_slide(self, slide, slide_num: int, temp_dir: str) -> str:
        """현재 슬라이드만 포함하는 임시 PPTX 파일 생성"""
        try:
            from pptx import Presentation
            
            # 새 프레젠테이션 생성
            prs = Presentation()
            
            # 현재 슬라이드를 새 프레젠테이션에 복사
            # (python-pptx의 제한으로 인해 복잡한 복사 로직이 필요할 수 있음)
            
            # 임시 파일 경로
            temp_pptx_path = os.path.join(temp_dir, f"slide_{slide_num}.pptx")
            
            # 간단한 방법: 빈 슬라이드만 생성하고 텍스트는 나중에 처리
            prs.save(temp_pptx_path)
            
            return temp_pptx_path
            
        except Exception as e:
            logger.warning(f"임시 PPTX 생성 실패: {e}")
            return None
    
    def _convert_pptx_to_pdf(self, pptx_path: str, temp_dir: str) -> str:
        """PPTX를 PDF로 변환"""
        try:
            from module.converters import ConverterFactory
            
            # LibreOffice 변환기 사용
            converter = ConverterFactory.get_best_converter("pptx", "pdf")
            if converter.is_available():
                pdf_path = os.path.join(temp_dir, f"slide_{Path(pptx_path).stem}.pdf")
                if converter.convert(pptx_path, pdf_path):
                    return pdf_path
            
            return None
            
        except Exception as e:
            logger.warning(f"PPTX -> PDF 변환 실패: {e}")
            return None
    
    def _convert_pdf_page_to_image(self, pdf_path: str, page_num: int, slide_num: int, temp_dir: str) -> str:
        """PDF 페이지를 이미지로 변환"""
        try:
            import fitz
            
            doc = fitz.open(pdf_path)
            if page_num <= len(doc):
                page = doc[page_num - 1]
                
                # 페이지를 이미지로 변환
                mat = fitz.Matrix(2.0, 2.0)  # 2배 확대
                pix = page.get_pixmap(matrix=mat)
                
                # 이미지 저장
                image_path = os.path.join(temp_dir, f"slide_{slide_num}.png")
                pix.save(image_path)
                
                doc.close()
                return image_path
            
            doc.close()
            return None
            
        except Exception as e:
            logger.warning(f"PDF -> 이미지 변환 실패: {e}")
            return None
    
    def _extract_pptx_image(self, shape, slide_num: int, element_index: int) -> str:
        """PPTX에서 이미지 추출"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            # 임시 디렉토리에 저장
            temp_dir = "temp_images"
            os.makedirs(temp_dir, exist_ok=True)
            
            image_filename = f"pptx_slide{slide_num}_img{element_index}.png"
            image_path = os.path.join(temp_dir, image_filename)
            
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            
            return image_path
            
        except Exception as e:
            raise Exception(f"이미지 추출 실패: {e}")
    
    def _process_xlsx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        XLSX 파일을 단순 변환 방식으로 처리
        
        요소 단위 분석만 수행 (Vision 분석 없음)
        text_chunk_to_embed에는 빈 문자열
        
        Args:
            file_path: XLSX 파일 경로
            
        Returns:
            요소 단위 분석 결과를 metadata에 포함한 청크 리스트
        """
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            all_elements = []
            
            logger.info(f"XLSX 단순 변환 처리 시작")
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # 시트별로 요소 분석
                table_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        table_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if table_data:
                    # 마크다운 테이블로 변환
                    table_text = self._table_to_markdown(table_data)
                    
                    element_info = {
                        "element_type": "table",
                        "content": table_text,
                        "table_format": "markdown",
                        "sheet_name": sheet_name,
                        "row_count": len(table_data),
                        "column_count": len(table_data[0]) if table_data else 0
                    }
                    all_elements.append(element_info)
            
            wb.close()
            
            # 결과 합성: text_chunk_to_embed는 빈 문자열, 요소 데이터는 metadata에
            simple_chunk = self._synthesize_simple_result(file_path, all_elements, "xlsx")
            
            logger.info(f"XLSX 단순 변환 처리 완료: {len(all_elements)}개 요소")
            return [simple_chunk]
            
        except Exception as e:
            logger.error(f"XLSX 파일 처리 오류: {e}")
            raise ContentExtractionError(f"XLSX 파일 처리 실패: {e}", file_path)
    
    def _process_pdf_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        PDF 파일을 이중 경로 하이브리드 아키텍처로 처리
        
        1. 요소 단위 분석: 텍스트, 테이블, 이미지 추출
        2. 페이지 단위 분석: 전체 페이지 Vision 분석
        
        Args:
            file_path: PDF 파일 경로
            
        Returns:
            이중 경로 분석 결과를 합성한 청크 리스트
        """
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            all_chunks = []
            
            logger.info(f"PDF 이중 경로 하이브리드 처리 시작: {len(doc)}개 페이지")
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                logger.info(f"페이지 {page_num + 1} 이중 경로 분석 중...")
                
                # 이중 경로 하이브리드 처리
                page_chunks = self._process_pdf_page_hybrid(page, page_num + 1, file_path)
                
                all_chunks.extend(page_chunks)
                logger.info(f"페이지 {page_num + 1} 이중 경로 처리 완료: {len(page_chunks)}개 청크")
            
            doc.close()
            logger.info(f"PDF 이중 경로 하이브리드 처리 완료: {len(all_chunks)}개 청크")
            return all_chunks
            
        except Exception as e:
            logger.error(f"PDF 파일 처리 오류: {e}")
            raise ContentExtractionError(f"PDF 파일 처리 실패: {e}", file_path)
    
    def _table_to_markdown(self, table_data: List[List[str]]) -> str:
        """테이블 데이터를 마크다운 형식으로 변환"""
        if not table_data:
            return ""
        
        lines = []
        
        # 헤더 행
        headers = table_data[0]
        lines.append("| " + " | ".join(str(cell) for cell in headers) + " |")
        lines.append("| " + " | ".join("---" for _ in headers) + " |")
        
        # 데이터 행들
        for row in table_data[1:]:
            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
        
        return "\n".join(lines)


class LayoutBasedProcessor(BaseProcessor):
    """Layout-based 파일 처리기 (PDF 변환 후 이미지 처리)"""
    
    def __init__(self, azure_processor, converter=None):
        super().__init__(azure_processor)
        self.converter = converter or ConverterFactory.get_best_converter("", "pdf")
    
    def process(self, file_path: str) -> List[Dict[str, Any]]:
        """Layout-based 파일을 PDF로 변환 후 처리하여 JSON 청크 리스트를 반환"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # PDF로 변환
                pdf_path = self._convert_to_pdf(file_path, temp_dir)
                
                # PDF를 이미지로 변환하여 처리
                return self._process_pdf_as_images(pdf_path, temp_dir)
                
        except Exception as e:
            logger.error(f"Layout-based 처리 오류: {e}")
            raise ProcessingError(f"Layout-based 처리 실패: {str(e)}", file_path)
    
    def _convert_to_pdf(self, file_path: str, temp_dir: str) -> str:
        """파일을 PDF로 변환"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.pdf':
                return file_path
            
            # 임시 PDF 파일 경로
            pdf_filename = f"converted_{Path(file_path).stem}.pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)
            
            # 변환기 선택 및 변환
            source_format = file_ext[1:]  # .docx -> docx
            converter = ConverterFactory.get_best_converter(source_format, "pdf")
            converter.convert(file_path, pdf_path)
            
            return pdf_path
            
        except Exception as e:
            logger.error(f"PDF 변환 오류: {e}")
            raise ProcessingError(f"PDF 변환 실패: {str(e)}", file_path)
    
    def _process_pdf_as_images(self, pdf_path: str, temp_dir: str) -> List[Dict[str, Any]]:
        """PDF를 이미지로 변환하여 처리"""
        try:
            if not self.azure_processor:
                logger.warning("Azure OpenAI 프로세서가 없어 텍스트 기반 처리로 대체")
                return self._process_pdf_text_based(pdf_path)
            
            doc = fitz.open(pdf_path)
            chunks = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 페이지를 이미지로 변환
                mat = fitz.Matrix(2.0, 2.0)  # 2배 확대
                pix = page.get_pixmap(matrix=mat)
                
                # 임시 이미지 저장
                image_filename = f"pdf_page{page_num + 1}.png"
                image_path = os.path.join(temp_dir, image_filename)
                pix.save(image_path)
                
                # GPT Vision으로 이미지 처리
                prompt = "이 페이지의 모든 텍스트 내용을 추출하고, 표나 이미지가 있다면 설명해주세요."
                text_content = self.azure_processor.image_to_text(image_path, prompt)
                
                # 결과를 청크로 저장
                metadata = {
                    "source_file": Path(pdf_path).name,
                    "section_title": f"페이지 {page_num + 1}",
                    "page_number": page_num + 1,
                    "element_type": "image_processed",
                    "processing_method": "gpt_vision",
                    "file_type": "pdf"
                }
                
                chunk = self._create_text_chunk(text_content, metadata)
                chunks.append(chunk)
            
            doc.close()
            logger.info(f"PDF 이미지 처리 완료: {len(chunks)}개 청크 생성")
            return chunks
            
        except Exception as e:
            logger.error(f"PDF 이미지 처리 오류: {e}")
            # 이미지 처리 실패 시 텍스트 기반 처리로 대체
            logger.info("이미지 처리 실패, 텍스트 기반 처리로 대체")
            return self._process_pdf_text_based(pdf_path)
    
    def _process_pdf_text_based(self, pdf_path: str) -> List[Dict[str, Any]]:
        """PDF를 텍스트 기반으로 처리 (이미지 처리 실패 시 대체 방법)"""
        try:
            doc = fitz.open(pdf_path)
            chunks = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                if text.strip():
                    metadata = {
                        "source_file": Path(pdf_path).name,
                        "section_title": f"페이지 {page_num + 1}",
                        "page_number": page_num + 1,
                        "element_type": "text",
                        "processing_method": "text_extraction",
                        "file_type": "pdf"
                    }
                    
                    chunk = self._create_text_chunk(text, metadata)
                    chunks.append(chunk)
            
            doc.close()
            logger.info(f"PDF 텍스트 기반 처리 완료: {len(chunks)}개 청크 생성")
            return chunks
            
        except Exception as e:
            logger.error(f"PDF 텍스트 기반 처리 오류: {e}")
            raise ContentExtractionError(f"PDF 텍스트 기반 처리 실패: {str(e)}", pdf_path)
    
    def _process_xml(self, file_path: str) -> List[Dict[str, Any]]:
        """
        XML 파일을 처리하여 JSON 청크 리스트를 반환
        
        Args:
            file_path: XML 파일 경로
            
        Returns:
            XML 파싱 결과를 포함한 청크 리스트
        """
        try:
            import xml.etree.ElementTree as ET
            import html
            
            logger.info(f"XML 파일 처리 시작: {file_path}")
            
            # XML 파싱
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            logger.info(f"XML 루트 요소: {root.tag}")
            
            chunks = []
            
            # RSS 형식인지 확인
            if root.tag == 'rss':
                chunks = self._process_rss_xml(root, file_path)
            else:
                # 일반 XML 처리
                chunks = self._process_generic_xml(root, file_path)
            
            if not chunks:
                logger.warning("처리된 요소가 없음, 전체 XML을 텍스트로 처리")
                chunks = self._process_xml_as_text(file_path)
            
            logger.info(f"XML 처리 완료: {len(chunks)}개 청크 생성")
            return chunks
            
        except ET.ParseError as e:
            logger.error(f"XML 파싱 오류: {e}")
            # 파싱 실패 시 텍스트로 처리
            return self._process_xml_as_text(file_path)
        except Exception as e:
            logger.error(f"XML 처리 오류: {e}")
            raise ContentExtractionError(f"XML 처리 실패: {str(e)}", file_path)
    
    def _process_rss_xml(self, root, file_path: str) -> List[Dict[str, Any]]:
        """RSS XML 처리 (JIRA RSS 특화)"""
        chunks = []
        
        # 채널 정보 추출
        channel = root.find('channel')
        if channel is not None:
            title_elem = channel.find('title')
            channel_title = title_elem.text if title_elem is not None else "Unknown"
            
            logger.info(f"RSS 채널: {channel_title}")
            
            # 이슈 정보 추출
            issue_elem = channel.find('issue')
            if issue_elem is not None:
                total_issues = issue_elem.get('total', '0')
                logger.info(f"총 이슈 수: {total_issues}")
            
            # 각 아이템(이슈) 처리
            items = channel.findall('item')
            logger.info(f"처리할 아이템 수: {len(items)}")
            
            for idx, item in enumerate(items):
                try:
                    # 이슈 정보 추출
                    issue_data = self._extract_jira_issue_data(item)
                    
                    if issue_data:
                        # 구조화된 텍스트 생성
                        structured_text = self._create_structured_issue_text(issue_data)
                        
                        # 청크 생성
                        metadata = {
                            "source": "jira_rss_item",
                            "issue_key": issue_data.get('key', ''),
                            "issue_type": issue_data.get('type', ''),
                            "status": issue_data.get('status', ''),
                            "priority": issue_data.get('priority', ''),
                            "assignee": issue_data.get('assignee', ''),
                            "reporter": issue_data.get('reporter', ''),
                            "created": issue_data.get('created', ''),
                            "updated": issue_data.get('updated', ''),
                            "item_index": idx + 1,
                            "file_type": "xml",
                            "processing_method": "rss_parser"
                        }
                        
                        chunk = self._create_text_chunk(structured_text, metadata)
                        chunks.append(chunk)
                        
                        if (idx + 1) % 100 == 0:
                            logger.info(f"{idx + 1}개 이슈 처리 완료...")
                
                except Exception as e:
                    logger.warning(f"아이템 {idx + 1} 처리 오류: {e}")
                    continue
        
        return chunks
    
    def _extract_jira_issue_data(self, item) -> Dict[str, str]:
        """JIRA 이슈 데이터 추출"""
        import html
        import re
        
        issue_data = {}
        
        try:
            # 기본 정보
            title_elem = item.find('title')
            if title_elem is not None:
                issue_data['title'] = title_elem.text or ''
            
            key_elem = item.find('key')
            if key_elem is not None:
                issue_data['key'] = key_elem.text or ''
            
            summary_elem = item.find('summary')
            if summary_elem is not None:
                issue_data['summary'] = summary_elem.text or ''
            
            # 설명 (HTML 디코딩)
            description_elem = item.find('description')
            if description_elem is not None and description_elem.text:
                # HTML 엔티티 디코딩
                description = html.unescape(description_elem.text)
                # HTML 태그 제거
                description = re.sub(r'<[^>]+>', '', description)
                issue_data['description'] = description.strip()
            
            # 프로젝트 정보
            project_elem = item.find('project')
            if project_elem is not None:
                issue_data['project'] = project_elem.text or ''
                issue_data['project_key'] = project_elem.get('key', '')
            
            # 이슈 타입
            type_elem = item.find('type')
            if type_elem is not None:
                issue_data['type'] = type_elem.text or ''
            
            # 상태
            status_elem = item.find('status')
            if status_elem is not None:
                issue_data['status'] = status_elem.text or ''
            
            # 우선순위
            priority_elem = item.find('priority')
            if priority_elem is not None:
                issue_data['priority'] = priority_elem.text or ''
            
            # 담당자
            assignee_elem = item.find('assignee')
            if assignee_elem is not None:
                issue_data['assignee'] = assignee_elem.text or ''
            
            # 보고자
            reporter_elem = item.find('reporter')
            if reporter_elem is not None:
                issue_data['reporter'] = reporter_elem.text or ''
            
            # 날짜 정보
            created_elem = item.find('created')
            if created_elem is not None:
                issue_data['created'] = created_elem.text or ''
            
            updated_elem = item.find('updated')
            if updated_elem is not None:
                issue_data['updated'] = updated_elem.text or ''
            
            resolved_elem = item.find('resolved')
            if resolved_elem is not None:
                issue_data['resolved'] = resolved_elem.text or ''
            
            # 링크
            link_elem = item.find('link')
            if link_elem is not None:
                issue_data['link'] = link_elem.text or ''
            
        except Exception as e:
            logger.warning(f"이슈 데이터 추출 오류: {e}")
        
        return issue_data
    
    def _create_structured_issue_text(self, issue_data: Dict[str, str]) -> str:
        """구조화된 이슈 텍스트 생성"""
        text_parts = []
        
        # 이슈 키와 제목
        if issue_data.get('key'):
            text_parts.append(f"이슈 키: {issue_data['key']}")
        if issue_data.get('title'):
            text_parts.append(f"제목: {issue_data['title']}")
        
        # 요약
        if issue_data.get('summary'):
            text_parts.append(f"요약: {issue_data['summary']}")
        
        # 설명
        if issue_data.get('description'):
            text_parts.append(f"설명: {issue_data['description']}")
        
        # 메타데이터
        metadata_parts = []
        if issue_data.get('type'):
            metadata_parts.append(f"타입: {issue_data['type']}")
        if issue_data.get('status'):
            metadata_parts.append(f"상태: {issue_data['status']}")
        if issue_data.get('priority'):
            metadata_parts.append(f"우선순위: {issue_data['priority']}")
        if issue_data.get('assignee'):
            metadata_parts.append(f"담당자: {issue_data['assignee']}")
        if issue_data.get('reporter'):
            metadata_parts.append(f"보고자: {issue_data['reporter']}")
        if issue_data.get('created'):
            metadata_parts.append(f"생성일: {issue_data['created']}")
        if issue_data.get('updated'):
            metadata_parts.append(f"수정일: {issue_data['updated']}")
        
        if metadata_parts:
            text_parts.append("메타데이터: " + ", ".join(metadata_parts))
        
        return "\n".join(text_parts)
    
    def _process_generic_xml(self, root, file_path: str) -> List[Dict[str, Any]]:
        """일반 XML 처리"""
        chunks = []
        
        def extract_text_recursive(element, depth: int = 0) -> str:
            """XML 요소에서 텍스트 재귀적으로 추출"""
            text_parts = []
            
            # 현재 요소의 텍스트
            if element.text and element.text.strip():
                text_parts.append(element.text.strip())
            
            # 자식 요소들 처리
            for child in element:
                child_text = extract_text_recursive(child, depth + 1)
                if child_text:
                    text_parts.append(child_text)
            
            # 현재 요소의 tail 텍스트
            if element.tail and element.tail.strip():
                text_parts.append(element.tail.strip())
            
            return " ".join(text_parts)
        
        # 루트 요소에서 텍스트 추출
        full_text = extract_text_recursive(root)
        
        if full_text.strip():
            metadata = {
                "source": "generic_xml",
                "root_element": root.tag,
                "file_type": "xml",
                "processing_method": "generic_parser"
            }
            
            chunk = self._create_text_chunk(full_text, metadata)
            chunks.append(chunk)
        
        return chunks
    
    def _process_xml_as_text(self, file_path: str) -> List[Dict[str, Any]]:
        """XML을 일반 텍스트로 처리"""
        try:
            import html
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # HTML 엔티티 디코딩
            content = html.unescape(content)
            
            metadata = {
                "source": "xml_as_text",
                "file_type": "xml",
                "processing_method": "text_fallback"
            }
            
            chunk = self._create_text_chunk(content, metadata)
            return [chunk]
            
        except Exception as e:
            logger.error(f"XML 텍스트 처리 오류: {e}")
            raise ContentExtractionError(f"XML 텍스트 처리 실패: {str(e)}", file_path) 